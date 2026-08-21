import zipfile
from pathlib import Path
from typing import Any

import pytest
from docx import Document

from app.api.source_routes import SUPPORTED_SOURCE_SUFFIXES, _validate_source_file
from app.extensions import default_extension_runtime
from app.services.parser_chain_execution import ParserChainExecution
from app.services.parsers import (
    MINERU_CAPABLE_SUFFIXES,
    MINERU_FALLBACK_WARNING_SUFFIXES,
    mineru_content_list_to_elements,
    mineru_label_prefix,
    parse_docx,
    parse_pptx,
    parse_source_file,
    parse_xls,
    parse_xlsx,
)

_FIXTURES_DIR = Path(__file__).parent / "fixtures"
# 生成方式（xlrd 只读不写，需要 xlwt 才能造 BIFF 文件；xlwt 不是仓库运行时依赖，
# 只临时装到隔离环境即可重造）：
#   python -m venv /tmp/xlwt-env && /tmp/xlwt-env/bin/pip install xlwt
#   /tmp/xlwt-env/bin/python - <<'PY'
#   import xlwt
#   book = xlwt.Workbook()
#   s1 = book.add_sheet("Sheet1")
#   for col, text in enumerate(["Name", "Count", "Score", "Joined", "Active"]):
#       s1.write(0, col, text)
#   date_style = xlwt.XFStyle()
#   date_style.num_format_str = "YYYY-MM-DD"
#   import datetime
#   s1.write(1, 0, "Alice"); s1.write(1, 1, 123); s1.write(1, 2, 1.5)
#   s1.write(1, 3, datetime.date(2024, 3, 15), date_style); s1.write(1, 4, True)
#   # row 2 left entirely blank on purpose (tests the skip-empty-row path)
#   s1.write(3, 0, "Bob"); s1.write(3, 1, 456)
#   s2 = book.add_sheet("Sheet2")
#   for row, text in enumerate(["Notes", "Second sheet row one.", "Second sheet row two."]):
#       s2.write(row, 0, text)
#   book.save("sample_legacy.xls")
#   PY
# datemode=0（1900 纪元，xlwt 默认）。
_LEGACY_XLS_FIXTURE = _FIXTURES_DIR / "sample_legacy.xls"


def _content_list():
    return [
        {"type": "title", "text": "Heading", "text_level": 1, "page_idx": 0},
        {"type": "text", "text": "Body text.", "page_idx": 0},
    ]


def test_mapper_default_prefix_is_pdf():
    els = mineru_content_list_to_elements("s1", _content_list())
    assert els[0].location_label.startswith("PDF p.1")
    assert all(e.metadata.get("parser") == "mineru" for e in els)
    assert els[0].metadata.get("source_format") == "pdf"


def test_mapper_custom_prefix_for_office():
    els = mineru_content_list_to_elements("s1", _content_list(), label_prefix="DOCX")
    assert els[0].location_label.startswith("DOCX p.1")
    assert els[0].metadata.get("source_format") == "docx"
    assert els[0].element_type == "heading"
    assert els[1].element_type == "paragraph"


def test_image_block_persists_asset_and_keeps_caption():
    content = [{"type": "image", "img_path": "images/fig1.jpg",
                "image_caption": ["Figure 1: layout"], "page_idx": 0}]
    saved = {}
    def persist(b, name):
        saved[name] = b
        return "asset-xyz"
    els = mineru_content_list_to_elements("s1", content, images={"fig1.jpg": b"J"}, persist_image=persist)
    img = [e for e in els if e.element_type == "image"]
    assert len(img) == 1
    assert img[0].metadata["asset_id"] == "asset-xyz"
    assert "Figure 1: layout" in img[0].text
    assert saved == {"fig1.jpg": b"J"}


def test_image_block_without_caption_is_not_dropped():
    content = [{"type": "image", "img_path": "images/x.png", "page_idx": 2}]
    els = mineru_content_list_to_elements("s1", content, images={"x.png": b"P"},
                                          persist_image=lambda b, n: "a1")
    img = [e for e in els if e.element_type == "image"]
    assert len(img) == 1                      # 旧行为会 continue 丢弃
    assert img[0].metadata["asset_id"] == "a1"
    assert img[0].text.strip()                # 占位文本非空


def test_image_block_no_persist_degrades_to_caption_text_only():
    content = [{"type": "image", "img_path": "images/x.png",
                "image_caption": ["cap"], "page_idx": 0}]
    els = mineru_content_list_to_elements("s1", content)   # 无 images/persist
    img = [e for e in els if e.element_type == "image"]
    assert len(img) == 1
    assert "asset_id" not in img[0].metadata
    assert "cap" in img[0].text


class FakeMineru:
    """模式无关的假 MinerU 客户端：按构造参数返回 content_list / 抛错 / 报未配置。"""

    def __init__(self, configured=True, content_list=None, raises=None, mode="http", images=None):
        self.configured = configured
        self._content_list = content_list if content_list is not None else []
        self._raises = raises
        self.mode = mode
        self._images = images or {}
        self.last_error = ""
        self.calls = []

    def parse(self, file_path, file_name):
        return self.parse_with_images(file_path, file_name)[0]

    def parse_with_images(self, file_path, file_name):
        # 故意不在抛出前预设 last_error：让生产代码的 "not last_error → 设置" 分支
        # 真正被执行（否则该分支被测试 fake 短路，删了也察觉不到）。
        self.last_error = ""
        self.calls.append((file_path, file_name))
        if self._raises is not None:
            raise self._raises
        return self._content_list, dict(self._images)


class _CloudOff:
    configured = False
    last_error = ""


class _NoConnection:
    @staticmethod
    def is_connection_held() -> bool:
        return False


def _parse_via_chain(
    source_id: str,
    path: Path,
    file_name: str,
    client: FakeMineru,
    persist_image=None,
):
    result = ParserChainExecution(
        host=default_extension_runtime().parser_chain,
        source_id=source_id,
        source_kind="file",
        file_path=str(path),
        file_name=file_name,
        source_url="",
        mineru_client=client,
        cloud_client=_CloudOff(),
        connection=_NoConnection(),
        make_persist_image=lambda: persist_image,
        delete_source_images=lambda: None,
        event_sink=lambda _event: None,
    ).run()
    client.last_error = result.mineru_error
    return list(result.elements)


def _make_docx(path: Path) -> Path:
    doc = Document()
    doc.add_paragraph("Hello from docx.")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "A"
    table.rows[0].cells[1].text = "B"
    doc.save(str(path))
    return path


def test_docx_uses_mineru_when_configured(tmp_path):
    path = _make_docx(tmp_path / "a.docx")
    client = FakeMineru(content_list=[{"type": "text", "text": "From MinerU.", "page_idx": 0}])
    els = _parse_via_chain("s1", path, "a.docx", client)
    assert any(e.metadata.get("parser") == "mineru" for e in els)
    assert els[0].location_label.startswith("DOCX p.1")


def test_docx_falls_back_when_not_configured(tmp_path):
    path = _make_docx(tmp_path / "a.docx")
    client = FakeMineru(configured=False)
    els = _parse_via_chain("s1", path, "a.docx", client)
    assert all(e.metadata.get("parser") == "mammoth" for e in els)
    assert any("Hello from docx." in e.text for e in els)


def test_docx_falls_back_on_mineru_error(tmp_path):
    path = _make_docx(tmp_path / "a.docx")
    client = FakeMineru(raises=RuntimeError("mineru boom"))
    els = _parse_via_chain("s1", path, "a.docx", client)  # 不应冒泡异常
    assert all(e.metadata.get("parser") == "mammoth" for e in els)
    assert client.last_error == "mineru boom"


def test_docx_falls_back_when_mineru_empty(tmp_path):
    path = _make_docx(tmp_path / "a.docx")
    client = FakeMineru(content_list=[])
    els = _parse_via_chain("s1", path, "a.docx", client)
    assert all(e.metadata.get("parser") == "mammoth" for e in els)


# -- DOCX 兜底：mammoth 语义化 -----------------------------------------------
#
# python-docx 逐段落取 `.text` 会把标题层级、列表和表格结构全部拍平；mammoth 把
# Word 样式映射成语义化 HTML，我们再折回仓库的元素词表。


def _make_rich_docx(path: Path) -> Path:
    doc = Document()
    doc.add_heading("Chapter One", level=1)
    doc.add_paragraph("Body para one.")
    doc.add_paragraph("bullet a", style="List Bullet")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "H1"
    table.rows[0].cells[1].text = "H2"
    table.rows[1].cells[0].text = "v1"
    table.rows[1].cells[1].text = "v2"
    doc.save(str(path))
    return path


def test_docx_fallback_keeps_heading_level_list_and_table(tmp_path):
    path = _make_rich_docx(tmp_path / "rich.docx")
    els = _parse_via_chain("s1", path, "rich.docx", FakeMineru(configured=False))
    assert all(e.metadata.get("parser") == "mammoth" for e in els)

    headings = [e for e in els if e.element_type == "heading"]
    assert [e.text for e in headings] == ["Chapter One"]
    assert headings[0].metadata["heading_level"] == 1
    assert headings[0].location_label == "DOCX heading 1"

    paragraphs = [e for e in els if e.element_type == "paragraph"]
    assert "Body para one." in [e.text for e in paragraphs]
    assert any(e.text.startswith("- ") and "bullet a" in e.text for e in paragraphs)

    tables = [e for e in els if e.element_type == "table"]
    assert len(tables) == 1                       # 旧兜底把表格拍成逐行 table_row
    assert "<table>" in tables[0].metadata["table_html"]
    assert "v1" in tables[0].metadata["table_html"]
    assert "H1 | H2" in tables[0].text


def _incompressible_png(width: int = 400, height: int = 400) -> bytes:
    """构造一张体积可观、几乎不可压缩的合法 PNG（约 480 KB）。

    小占位图证明不了任何事：mammoth 默认 handler 会把**每一次**图片引用 base64
    物化进 HTML，只有真实体积才能让「物化 vs 不物化」在断言上分得开。
    """
    import os
    import struct
    import zlib

    raw = b"".join(b"\x00" + os.urandom(width * 3) for _ in range(height))

    def _chunk(tag: bytes, data: bytes) -> bytes:
        body = tag + data
        return struct.pack(">I", len(data)) + body + struct.pack(">I", zlib.crc32(body) & 0xFFFFFFFF)

    return (
        b"\x89PNG\r\n\x1a\n"
        + _chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0))
        + _chunk(b"IDAT", zlib.compress(raw, 0))
        + _chunk(b"IEND", b"")
    )


def test_docx_fallback_never_materialises_inline_image_bytes(tmp_path, monkeypatch):
    """mammoth 默认 handler 把每张内嵌图 base64 物化进 HTML，而本解析器 100% 丢弃
    这些字节：一张 480 KB 的图引用 6 次就是 ~4 MB HTML + 峰值内存，全是白付的。

    断言钉的是 **convert_image 接线真的生效**——包住真实 mammoth 转换、量它产出的
    HTML。删掉 convert_image 参数会让 HTML 从几 KB 涨到 MB 级并出现 data:image，
    这条用例立刻变红（只断言元素里没有 base64 是空守卫：元素本来就丢 <img>）。
    """
    import mammoth

    image_path = tmp_path / "big.png"
    image_path.write_bytes(_incompressible_png())
    doc = Document()
    doc.add_paragraph("before")
    for _ in range(6):
        doc.add_picture(str(image_path))
    doc.add_paragraph("after")
    path = tmp_path / "img.docx"
    doc.save(str(path))

    seen = {}
    real_convert = mammoth.convert_to_html

    def _spy(*args, **kwargs):
        result = real_convert(*args, **kwargs)
        seen["markup"] = str(result.value or "")
        return result

    monkeypatch.setattr(mammoth, "convert_to_html", _spy)
    els = _parse_via_chain("s1", path, "img.docx", FakeMineru(configured=False))

    markup = seen["markup"]
    assert "data:image" not in markup                 # 默认 handler 会有 6 处
    assert len(markup) < 10_000                       # 默认 handler ≈ 4 MB
    assert image_path.stat().st_size > 400_000        # 反向自证：图确实很大
    assert [e.text for e in els] == ["before", "after"]
    assert not any("base64" in str(e.metadata) or "base64" in e.text for e in els)


def test_docx_fallback_unescapes_entities_in_table_text(tmp_path):
    """表格 text 进 embedding/FTS/展示，不能带 HTML 实体：`R&D` 曾变成 `R&amp;D`。
    table_html metadata 仍保持转义（渲染端要的是 HTML）。"""
    doc = Document()
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "R&D"
    table.rows[0].cells[1].text = "<b>literal</b>"
    path = tmp_path / "amp.docx"
    doc.save(str(path))

    els = _parse_via_chain("s1", path, "amp.docx", FakeMineru(configured=False))
    tables = [e for e in els if e.element_type == "table"]
    assert len(tables) == 1
    assert "R&D" in tables[0].text
    assert "&amp;" not in tables[0].text
    assert "<b>literal</b>" in tables[0].text          # 反转义在剥标签之后才安全
    assert "R&amp;D" in tables[0].metadata["table_html"]


def test_docx_fallback_splits_long_tables_into_bounded_elements(tmp_path):
    """`chunking.py` 从不切分单个元素：一张大表折成一个元素 = 一个几十万字符的
    chunk（embedding 截断 + 召回坍缩）。按行切成多个可检索单位，且一行不丢。"""
    rows = 65
    doc = Document()
    table = doc.add_table(rows=rows, cols=2)
    for index in range(rows):
        table.rows[index].cells[0].text = f"r{index}c0"
        table.rows[index].cells[1].text = f"r{index}c1"
    path = tmp_path / "long.docx"
    doc.save(str(path))

    els = _parse_via_chain("s1", path, "long.docx", FakeMineru(configured=False))
    tables = [e for e in els if e.element_type == "table"]
    assert len(tables) == 3                                   # 65 行 / 每段 30 行
    assert [e.location_label for e in tables] == [
        "DOCX table 1 part 1",
        "DOCX table 1 part 2",
        "DOCX table 1 part 3",
    ]
    assert [e.metadata["table_part"] for e in tables] == [1, 2, 3]
    assert all(e.metadata["table_parts"] == 3 for e in tables)
    assert all(e.metadata["table_index"] == 1 for e in tables)
    assert all(e.metadata["table_html"].startswith("<table>") for e in tables)
    assert all(e.metadata["table_html"].endswith("</table>") for e in tables)
    assert max(len(e.text) for e in tables) < 2_000            # 单元素有界

    merged = " ; ".join(e.text for e in tables)
    for index in range(rows):                                  # 合并后一行不丢
        assert f"r{index}c0 | r{index}c1" in merged
    assert sum(e.metadata["table_html"].count("<tr>") for e in tables) == rows


def test_split_table_html_defers_the_split_while_a_rowspan_is_active():
    """codex R2 P2：rowspan 跨段会让后段整列缺失、其余单元格全部错位。切点只落在
    没有活跃竖向合并的行边界；跨越默认切点（第 30 行）的合并把收段顺延到 span
    闭合之后，被合并覆盖的行必须与 span 起始行同段。"""
    from app.services.parsers import _split_table_html

    rows = []
    for index in range(65):
        if index == 28:
            # 覆盖 28..31 四行的竖向合并，恰好跨过 30 行的默认切点。
            rows.append(f'<tr><td rowspan="4">merged</td><td>r{index}</td></tr>')
        elif 29 <= index <= 31:
            rows.append(f"<tr><td>r{index}</td></tr>")
        else:
            rows.append(f"<tr><td>a{index}</td><td>r{index}</td></tr>")
    raw = "<table>" + "".join(rows) + "</table>"

    parts = _split_table_html(raw)
    assert len(parts) == 3
    assert all(p.startswith("<table>") and p.endswith("</table>") for p in parts)
    # 首段顺延到 span 闭合（32 行），span 起始行与被覆盖行都在首段。
    assert parts[0].count("<tr") == 32
    assert 'rowspan="4"' in parts[0]
    for covered in ("<td>r29</td>", "<td>r30</td>", "<td>r31</td>"):
        assert covered in parts[0]
        assert all(covered not in later for later in parts[1:])
    assert sum(p.count("<tr") for p in parts) == 65  # 顺延不丢行


def test_docx_fallback_keeps_small_tables_as_one_element(tmp_path):
    """未超限的表格逐字不变：不加 part 后缀、不加 part metadata。"""
    path = _make_rich_docx(tmp_path / "small.docx")
    els = _parse_via_chain("s1", path, "small.docx", FakeMineru(configured=False))
    tables = [e for e in els if e.element_type == "table"]
    assert len(tables) == 1
    assert tables[0].location_label == "DOCX table 1"
    assert "table_part" not in tables[0].metadata


def test_mammoth_block_parser_recovers_an_unclosed_table():
    """防御性：</table> 缺失时不能把整张表连同其后的正文一起吞掉。"""
    from app.services.parsers import _mammoth_html_blocks

    blocks = _mammoth_html_blocks("<p>lead</p><table><tr><td>a</td><td>b</td></tr>")
    assert [b[0] for b in blocks] == ["paragraph", "table"]
    assert blocks[1][2] == "a | b"


def test_docx_falls_back_to_python_docx_when_mammoth_fails(tmp_path, monkeypatch):
    """mammoth 缺失/出错是兜底链的最后一环触发条件，不能让摄取整体失败。"""
    import mammoth

    path = _make_rich_docx(tmp_path / "rich.docx")

    def _boom(*args, **kwargs):
        raise RuntimeError("mammoth boom")

    monkeypatch.setattr(mammoth, "convert_to_html", _boom)
    els = _parse_via_chain("s1", path, "rich.docx", FakeMineru(configured=False))
    assert els
    assert all(e.metadata.get("parser") == "docx" for e in els)
    assert any("Body para one." in e.text for e in els)


def test_parse_source_file_forwards_client_to_docx(tmp_path):
    path = _make_docx(tmp_path / "a.docx")
    client = FakeMineru(content_list=[{"type": "text", "text": "From MinerU.", "page_idx": 0}])
    els = _parse_via_chain("s1", path, "a.docx", client)
    assert any(e.metadata.get("parser") == "mineru" for e in els)


_SLIDE_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    "<p:cSld><p:spTree><p:sp><p:txBody>"
    "<a:p><a:r><a:t>{text}</a:t></a:r></a:p>"
    "</p:txBody></p:sp></p:spTree></p:cSld></p:sld>"
)


def _make_pptx(path: Path, slide_texts) -> Path:
    with zipfile.ZipFile(path, "w") as archive:
        for index, text in enumerate(slide_texts, start=1):
            archive.writestr(f"ppt/slides/slide{index}.xml", _SLIDE_XML.format(text=text))
    return path


def test_pptx_uses_mineru_when_configured(tmp_path):
    path = _make_pptx(tmp_path / "a.pptx", ["Slide one body"])
    client = FakeMineru(content_list=[{"type": "text", "text": "From MinerU.", "page_idx": 0}])
    els = _parse_via_chain("s1", path, "a.pptx", client)
    assert any(e.metadata.get("parser") == "mineru" for e in els)
    assert els[0].location_label.startswith("PPTX p.1")


def test_pptx_falls_back_when_not_configured(tmp_path):
    path = _make_pptx(tmp_path / "a.pptx", ["Slide one body"])
    client = FakeMineru(configured=False)
    els = _parse_via_chain("s1", path, "a.pptx", client)
    assert all(e.metadata.get("parser") == "pptx" for e in els)
    assert any("Slide one body" in e.text for e in els)


def test_pptx_falls_back_on_mineru_error(tmp_path):
    path = _make_pptx(tmp_path / "a.pptx", ["Slide one body"])
    client = FakeMineru(raises=RuntimeError("mineru boom"))
    els = _parse_via_chain("s1", path, "a.pptx", client)
    assert all(e.metadata.get("parser") == "pptx" for e in els)
    assert client.last_error == "mineru boom"


def test_pptx_falls_back_when_mineru_empty(tmp_path):
    path = _make_pptx(tmp_path / "a.pptx", ["Slide one body"])
    client = FakeMineru(content_list=[])
    els = _parse_via_chain("s1", path, "a.pptx", client)
    assert all(e.metadata.get("parser") == "pptx" for e in els)


# -- PPTX 兜底：python-pptx ---------------------------------------------------
#
# 原始 XML 抽取只走 `p:sp` 形状，幻灯片里的表格（`p:graphicFrame`/`a:tbl`）与图表
# 标题**整块丢失**。下面的表格断言就是这条特性的核心回归钉。


def _make_real_pptx(path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Slide body text"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(1, 0).text = "v1"
    table.cell(1, 1).text = "v2"
    slide.notes_slide.notes_text_frame.text = "Speaker note here"
    presentation.save(str(path))
    return path


def test_pptx_fallback_recovers_slide_tables_and_notes(tmp_path):
    path = _make_real_pptx(tmp_path / "real.pptx")
    els = _parse_via_chain("s1", path, "real.pptx", FakeMineru(configured=False))
    assert all(e.metadata.get("parser") == "python-pptx" for e in els)

    texts = [e for e in els if e.element_type == "slide_text"]
    assert any("Slide body text" in e.text for e in texts)
    assert texts[0].metadata["slide_number"] == 1

    tables = [e for e in els if e.element_type == "table"]
    assert len(tables) == 1                       # 旧兜底 (parse_pptx_basic) 会丢掉它
    assert "H1 | H2" in tables[0].text
    assert "<table>" in tables[0].metadata["table_html"]
    assert tables[0].location_label == "PPTX slide 1 table 1"

    notes = [e for e in els if e.element_type == "speaker_notes"]
    assert [e.text for e in notes] == ["Speaker note here"]


def test_pptx_fallback_reads_grouped_shape_text(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    first.text_frame.text = "Grouped one"
    second = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
    second.text_frame.text = "Grouped two"
    slide.shapes.add_group_shape([first, second])
    path = tmp_path / "group.pptx"
    presentation.save(str(path))

    els = _parse_via_chain("s1", path, "group.pptx", FakeMineru(configured=False))
    assert {"Grouped one", "Grouped two"} <= {e.text for e in els}


def test_pptx_fallback_reads_nested_group_shape_text(tmp_path):
    """组合形状要**完全展开且无深度上限**（codex R3 P2）：任何深度截断都是对用户
    数据的静默丢弃，还不如它取代的原始 XML 抽取（`.//p:sp` 本来就是任意深度）。
    20 层嵌套钉住「不存在护栏深度常量」——旧实现的 16 层上限在这里必红。"""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape([])
    for _ in range(19):
        group = group.shapes.add_group_shape([])
    box = group.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    box.text_frame.text = "Depth twenty text"
    shallow = slide.shapes.add_group_shape([])
    shallow_box = shallow.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
    shallow_box.text_frame.text = "Depth one text"
    path = tmp_path / "nested.pptx"
    presentation.save(str(path))

    els = _parse_via_chain("s1", path, "nested.pptx", FakeMineru(configured=False))
    assert all(e.metadata.get("parser") == "python-pptx" for e in els)
    texts = {e.text for e in els}
    assert "Depth twenty text" in texts
    assert "Depth one text" in texts


def test_pptx_fallback_survives_a_shape_with_unimplemented_shape_type(tmp_path):
    """python-pptx 对没有可识别几何的 `p:sp` 在 `.shape_type` 里 **raise
    NotImplementedError**（getattr 的 default 挡不住）。判组读它 = 一个畸形形状让
    整份 deck 掉回 parse_pptx_basic、幻灯片表格全丢。"""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    broken = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    broken.text_frame.text = "Broken shape text"
    # 去掉 txBox 标记与 prstGeom：既不是占位符、不是 autoshape、不是 textbox、
    # 也没有 custGeom —— 正是 shape_type 抛 NotImplementedError 的形状。
    sp = broken._element
    sp.nvSpPr.cNvSpPr.attrib.pop("txBox", None)
    prst = sp.spPr.prstGeom
    sp.spPr.remove(prst)
    good = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
    good.text_frame.text = "Healthy shape text"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    path = tmp_path / "broken.pptx"
    presentation.save(str(path))

    # 自证前提：这个形状的 shape_type 确实会抛。
    with pytest.raises(NotImplementedError):
        _ = broken.shape_type

    els = _parse_via_chain("s1", path, "broken.pptx", FakeMineru(configured=False))
    assert all(e.metadata.get("parser") == "python-pptx" for e in els)   # 未掉回原始 XML
    assert "Healthy shape text" in {e.text for e in els}
    assert any(e.element_type == "table" for e in els)                   # 表格没被连坐


def test_pptx_fallback_isolates_a_shape_that_raises_while_being_read(tmp_path, monkeypatch):
    """per-shape 兜底：坏形状只丢自己。没有它，任一形状在属性访问处抛异常就会由整
    函数的 fail-open 接住、整份 deck 掉回 parse_pptx_basic（丢表格、丢备注）。"""
    from app.services import parsers as parsers_module

    path = _make_real_pptx(tmp_path / "real.pptx")

    def _boom(table):
        raise RuntimeError("shape boom")

    monkeypatch.setattr(parsers_module, "_pptx_table_html", _boom)
    els = _parse_via_chain("s1", path, "real.pptx", FakeMineru(configured=False))
    assert all(e.metadata.get("parser") == "python-pptx" for e in els)   # 未掉回原始 XML
    assert not any(e.element_type == "table" for e in els)               # 坏的那个丢了
    assert any("Slide body text" in e.text for e in els)                 # 同页其他形状还在
    assert any(e.element_type == "speaker_notes" for e in els)           # 备注也没被连坐


def test_pptx_fallback_keeps_merged_cells_as_spans(tmp_path):
    """合并单元格：被覆盖的格跳过，起始格写 rowspan/colspan——否则跨 3 列的表头会
    被拍成 3 个格、其中 2 个空。"""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(6), Inches(1)).table
    table.cell(0, 0).merge(table.cell(0, 2))
    table.cell(0, 0).text = "Spanning header"
    table.cell(1, 0).text = "v0"
    table.cell(1, 1).text = "v1"
    table.cell(1, 2).text = "v2"
    path = tmp_path / "merged.pptx"
    presentation.save(str(path))

    els = _parse_via_chain("s1", path, "merged.pptx", FakeMineru(configured=False))
    tables = [e for e in els if e.element_type == "table"]
    assert len(tables) == 1
    html = tables[0].metadata["table_html"]
    assert 'colspan="3"' in html
    assert html.count("<td") == 4                       # 1 个合并头 + 3 个数据格
    assert "Spanning header" in tables[0].text
    assert "v0 | v1 | v2" in tables[0].text


def test_pptx_fallback_splits_long_tables_into_bounded_elements(tmp_path):
    """与 DOCX 同一条检索粒度边界：大表按行切成多个 table 元素，一行不丢。"""
    from pptx import Presentation
    from pptx.util import Inches

    rows = 65
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(rows, 2, Inches(1), Inches(1), Inches(4), Inches(1)).table
    for index in range(rows):
        table.cell(index, 0).text = f"r{index}c0"
        table.cell(index, 1).text = f"r{index}c1"
    path = tmp_path / "longtable.pptx"
    presentation.save(str(path))

    els = _parse_via_chain("s1", path, "longtable.pptx", FakeMineru(configured=False))
    tables = [e for e in els if e.element_type == "table"]
    assert len(tables) == 3
    assert [e.location_label for e in tables] == [
        "PPTX slide 1 table 1 part 1",
        "PPTX slide 1 table 1 part 2",
        "PPTX slide 1 table 1 part 3",
    ]
    assert all(e.metadata["table_index"] == 1 for e in tables)
    assert max(len(e.text) for e in tables) < 2_000
    merged = " ; ".join(e.text for e in tables)
    for index in range(rows):
        assert f"r{index}c0 | r{index}c1" in merged


def test_pptx_fallback_survives_missing_notes_text_frame(tmp_path, monkeypatch):
    """`notes_text_frame` 官方允许为 None（备注母版无 body 占位符/占位符被删而
    notesSlide part 仍在）；直接 `.text` 会 AttributeError 让整份 deck 掉回原始 XML。"""
    path = _make_real_pptx(tmp_path / "real.pptx")

    from pptx.slide import NotesSlide

    monkeypatch.setattr(
        NotesSlide, "notes_text_frame", property(lambda self: None), raising=True
    )
    els = _parse_via_chain("s1", path, "real.pptx", FakeMineru(configured=False))
    assert all(e.metadata.get("parser") == "python-pptx" for e in els)   # 未掉回原始 XML
    assert not any(e.element_type == "speaker_notes" for e in els)
    assert any("Slide body text" in e.text for e in els)


def test_pptx_falls_back_to_raw_xml_when_python_pptx_fails(tmp_path, monkeypatch):
    """python-pptx 缺失/包读不了时仍要出内容，原始 XML 抽取是最后一环。"""
    import pptx

    path = _make_real_pptx(tmp_path / "real.pptx")

    def _boom(*args, **kwargs):
        raise RuntimeError("pptx boom")

    monkeypatch.setattr(pptx, "Presentation", _boom)
    els = _parse_via_chain("s1", path, "real.pptx", FakeMineru(configured=False))
    assert els
    assert all(e.metadata.get("parser") == "pptx" for e in els)
    assert any("Slide body text" in e.text for e in els)


def test_parse_source_file_forwards_client_to_pptx(tmp_path):
    path = _make_pptx(tmp_path / "a.pptx", ["Slide one body"])
    client = FakeMineru(content_list=[{"type": "text", "text": "From MinerU.", "page_idx": 0}])
    els = _parse_via_chain("s1", path, "a.pptx", client)
    assert any(e.metadata.get("parser") == "mineru" for e in els)


def _make_xlsx(path: Path) -> Path:
    from openpyxl import Workbook

    book = Workbook()
    sheet = book.active
    sheet.title = "Sheet1"
    sheet.append(["Header A", "Header B"])
    sheet.append(["Hello from xlsx.", "42"])
    book.save(str(path))
    return path


def _make_xlsx_rows(path: Path, rows: int, cols: int = 2) -> Path:
    """行数与列数可控的工作簿——对账的两个分母分别是 rows 与 rows*cols。"""
    from openpyxl import Workbook

    book = Workbook()
    sheet = book.active
    sheet.title = "Sheet1"
    for index in range(rows):
        sheet.append([f"r{index}c{col}" for col in range(cols)])
    book.save(str(path))
    return path


def _mineru_table(tr_count: int, td_per_row: int = 2) -> list[dict]:
    rows = "".join(
        "<tr>" + "".join(f"<td>c{i}_{j}</td>" for j in range(td_per_row)) + "</tr>"
        for i in range(tr_count)
    )
    return [
        {
            "type": "table",
            "table_body": f"<table>{rows}</table>",
            "table_caption": ["Sheet1"],
            "page_idx": 0,
        }
    ]


def test_xlsx_uses_mineru_when_configured(tmp_path):
    """MinerU 原生解析 XLSX：表格保留 HTML 结构，而不是 openpyxl 的逐行拍平。"""
    path = _make_xlsx(tmp_path / "a.xlsx")
    client = FakeMineru(
        content_list=[
            {
                "type": "table",
                "table_body": (
                    "<table><tr><td>Header A</td><td>Header B</td></tr>"
                    "<tr><td>Hello from xlsx.</td><td>42</td></tr></table>"
                ),
                "table_caption": ["Sheet1"],
                "page_idx": 0,
            }
        ]
    )
    els = _parse_via_chain("s1", path, "a.xlsx", client)
    assert all(e.metadata.get("parser") == "mineru" for e in els)
    assert not any(e.metadata.get("parser") == "xlsx" for e in els)  # openpyxl 未被调用
    assert els[0].location_label.startswith("XLSX p.1")
    assert els[0].metadata.get("source_format") == "xlsx"
    assert "<table>" in els[0].metadata["table_html"]


# -- MinerU 工作簿产出的行+格覆盖对账 -----------------------------------------
#
# MinerU 按渲染页出表，打印区域/列宽/隐藏列会让它整列整 sheet 丢内容**且不报错**。
# 只要非空就采信 = 静默数据损失（parser 记成 mineru，连降级警告都不亮）。
# 只数行同样不够：宽表下 `<tr>` 常常足额而页宽之外的列被整片丢掉，行分子看不出来。


def test_xlsx_rejects_mineru_output_that_drops_rows(tmp_path):
    path = _make_xlsx_rows(tmp_path / "big.xlsx", 10)      # 10 行 × 2 列 = 20 格
    client = FakeMineru(content_list=_mineru_table(7))     # 7/10 < ceil(10*0.8)=8
    els = _parse_via_chain("s1", path, "big.xlsx", client)
    assert all(e.metadata.get("parser") == "xlsx" for e in els)   # 整份丢弃，回落 openpyxl
    assert len(els) == 10                                          # 全保真：一行不少
    assert client.last_error == (
        "MinerU workbook output covered 7/10 rows and 14/20 cells; using openpyxl"
    )


def test_xlsx_rejects_mineru_output_that_drops_columns(tmp_path):
    """行数足额、列被丢掉——只数 `<tr>` 的对账对这种宽表残缺完全瞎。

    MinerU 保住了每一行，但页宽之外的 4 列没了：行覆盖 10/10 满分，格覆盖只有
    10/50。这正是这份产出必须被拒的理由。
    """
    path = _make_xlsx_rows(tmp_path / "wide.xlsx", 10, cols=5)    # 10 行 × 5 列 = 50 格
    client = FakeMineru(content_list=_mineru_table(10, td_per_row=1))
    els = _parse_via_chain("s1", path, "wide.xlsx", client)
    assert all(e.metadata.get("parser") == "xlsx" for e in els)
    assert len(els) == 10
    assert client.last_error == (
        "MinerU workbook output covered 10/10 rows and 10/50 cells; using openpyxl"
    )


def test_xlsx_rejects_a_structurally_complete_but_value_empty_grid(tmp_path):
    """codex R6 P1：识别失败时 MinerU 常返回**结构完整、值却全空**的格子。按
    `<tr>/<td>` 标签计数的分子会被 10 行 20 个空 `<td>` 刷满两个阈值；分子必须与
    openpyxl 分母同一套非空语义——剥标签+反转义后有实文的格才算覆盖，空格不算，
    全空行也不进行分子。"""
    path = _make_xlsx_rows(tmp_path / "hollow.xlsx", 10)          # 10 行 × 2 列 = 20 格
    rows_html = "".join(
        "<tr>" + ("<td>v0</td><td>v1</td>" if i == 0 else "<td></td><td> &nbsp; </td>") + "</tr>"
        for i in range(10)
    )
    client = FakeMineru(content_list=[{
        "type": "table",
        "table_body": f"<table>{rows_html}</table>",
        "table_caption": ["Sheet1"],
        "page_idx": 0,
    }])
    els = _parse_via_chain("s1", path, "hollow.xlsx", client)
    assert all(e.metadata.get("parser") == "xlsx" for e in els)   # 拒收，回落 openpyxl
    assert len(els) == 10
    assert client.last_error == (
        "MinerU workbook output covered 1/10 rows and 2/20 cells; using openpyxl"
    )


def test_xlsx_accepts_a_wide_table_that_keeps_its_columns(tmp_path):
    """反向护栏：同一张宽表，列没丢就照常采信——格指标不能变成第二把误杀刀。"""
    path = _make_xlsx_rows(tmp_path / "wide.xlsx", 10, cols=5)
    client = FakeMineru(content_list=_mineru_table(10, td_per_row=5))  # 50/50 格
    els = _parse_via_chain("s1", path, "wide.xlsx", client)
    assert all(e.metadata.get("parser") == "mineru" for e in els)
    assert client.last_error == ""


def test_xlsx_accepts_mineru_output_at_the_coverage_threshold(tmp_path):
    path = _make_xlsx_rows(tmp_path / "big.xlsx", 10)
    # 8/10 行 == ceil(10*0.8)，16/20 格 == ceil(20*0.8)：两个维度都恰好压线。
    client = FakeMineru(content_list=_mineru_table(8))
    els = _parse_via_chain("s1", path, "big.xlsx", client)
    assert all(e.metadata.get("parser") == "mineru" for e in els)
    assert client.last_error == ""


def test_xlsx_coverage_counts_non_table_elements(tmp_path):
    """MinerU 会把 sheet 标题/游离单元格映射成标题或段落而不是表格行；不计入它们
    会把正常产出误判成缺行、把 MinerU 路径彻底废掉。

    双指标口径下，一条文本元素只向「行」与「格」各计 1——所以这里用单列工作簿，
    分母是 3 行 3 格，虚增无从发生。"""
    path = _make_xlsx_rows(tmp_path / "small.xlsx", 3, cols=1)
    client = FakeMineru(
        content_list=[
            {"type": "text", "text": f"cell {index}", "page_idx": 0} for index in range(3)
        ]
    )
    els = _parse_via_chain("s1", path, "small.xlsx", client)
    assert all(e.metadata.get("parser") == "mineru" for e in els)


def test_xlsx_coverage_ignores_image_elements(tmp_path):
    """图片不对应任何物理行或格，两个分子都不能计入它。

    这个形状此前会通过：4 张图各计 1 行 + 1 条文本 = 5/5 行覆盖满分，而工作簿里
    真正被 MinerU 取到的内容只有 1 行。
    """
    path = _make_xlsx_rows(tmp_path / "imgs.xlsx", 5, cols=1)     # 5 行 5 格
    client = FakeMineru(
        content_list=[
            {"type": "image", "img_path": f"images/p{i}.png", "page_idx": 0} for i in range(4)
        ]
        + [{"type": "text", "text": "Sheet1", "page_idx": 0}]
    )
    els = _parse_via_chain("s1", path, "imgs.xlsx", client)
    assert all(e.metadata.get("parser") == "xlsx" for e in els)   # 拒收
    assert client.last_error == (
        "MinerU workbook output covered 1/5 rows and 1/5 cells; using openpyxl"
    )


def test_xlsx_rejected_output_persists_no_images(tmp_path):
    """拒收路径必须一张图都没写：映射分两段（探针不带 persist_image，通过才重映射），
    否则被丢弃的产出会在资产表/磁盘上留下永远无人引用的孤儿图片。"""
    path = _make_xlsx_rows(tmp_path / "big.xlsx", 10)
    client = FakeMineru(
        content_list=_mineru_table(2)
        + [{"type": "image", "img_path": "images/a.png", "page_idx": 0}],
        images={"a.png": b"PNGDATA"},
    )
    calls = []

    def _persist(data, name):
        calls.append(name)
        return "asset-1"

    els = _parse_via_chain("s1", path, "big.xlsx", client, persist_image=_persist)
    assert all(e.metadata.get("parser") == "xlsx" for e in els)   # 拒收
    assert calls == []                                            # 零持久化


def test_xlsx_accepted_output_persists_images(tmp_path):
    """反向护栏：对账通过时图片照常持久化——两段式映射不能把图片一起弄丢。"""
    path = _make_xlsx_rows(tmp_path / "big.xlsx", 10)
    client = FakeMineru(
        content_list=_mineru_table(8)
        + [{"type": "image", "img_path": "images/a.png", "page_idx": 0}],
        images={"a.png": b"PNGDATA"},
    )
    calls = []

    def _persist(data, name):
        calls.append(name)
        return "asset-1"

    els = _parse_via_chain("s1", path, "big.xlsx", client, persist_image=_persist)
    assert all(e.metadata.get("parser") == "mineru" for e in els)
    assert calls == ["a.png"]
    assert any(e.metadata.get("asset_id") == "asset-1" for e in els)


def test_xlsx_coverage_is_fail_open_when_the_workbook_cannot_be_read(tmp_path):
    """对账是护栏，不能自己变成新的失败模式：openpyxl 读不了就采信 MinerU。"""
    path = tmp_path / "stub.xlsx"
    path.write_bytes(b"PK\x03\x04not-a-real-workbook")
    client = FakeMineru(content_list=[{"type": "text", "text": "From MinerU.", "page_idx": 0}])
    els = _parse_via_chain("s1", path, "stub.xlsx", client)
    assert all(e.metadata.get("parser") == "mineru" for e in els)
    assert client.last_error == ""


def test_xlsx_falls_back_when_not_configured(tmp_path):
    path = _make_xlsx(tmp_path / "a.xlsx")
    client = FakeMineru(configured=False)
    els = _parse_via_chain("s1", path, "a.xlsx", client)
    assert all(e.metadata.get("parser") == "xlsx" for e in els)
    assert any("Hello from xlsx." in e.text for e in els)


def test_xlsx_falls_back_on_mineru_error(tmp_path):
    path = _make_xlsx(tmp_path / "a.xlsx")
    client = FakeMineru(raises=RuntimeError("mineru boom"))
    els = _parse_via_chain("s1", path, "a.xlsx", client)  # 不应冒泡异常
    assert all(e.metadata.get("parser") == "xlsx" for e in els)
    assert client.last_error == "mineru boom"


def test_xlsx_falls_back_when_mineru_empty(tmp_path):
    path = _make_xlsx(tmp_path / "a.xlsx")
    client = FakeMineru(content_list=[])
    els = _parse_via_chain("s1", path, "a.xlsx", client)
    assert all(e.metadata.get("parser") == "xlsx" for e in els)
    assert client.last_error == "MinerU content_list mapped to zero source elements"


def test_parse_source_file_forwards_client_to_xlsx(tmp_path):
    path = _make_xlsx(tmp_path / "a.xlsx")
    client = FakeMineru(content_list=_mineru_table(2))
    els = _parse_via_chain("s1", path, "a.xlsx", client)
    assert any(e.metadata.get("parser") == "mineru" for e in els)


def test_parse_source_file_forwards_client_to_xlsm(tmp_path):
    # .xlsm 与 .xlsx 同为 OOXML 工作簿，同样要拿到 client（否则宏工作簿静默退化）。
    path = _make_xlsx(tmp_path / "a.xlsm")
    client = FakeMineru(content_list=_mineru_table(2))
    els = _parse_via_chain("s1", path, "a.xlsm", client)
    assert any(e.metadata.get("parser") == "mineru" for e in els)


# -- 旧版二进制 .xls（xlrd，无 MinerU 分支）------------------------------------


def test_parse_xls_reads_both_sheets_and_normalizes_cell_types():
    els = parse_xls("s1", _LEGACY_XLS_FIXTURE)
    assert all(e.element_type == "table_row" for e in els)
    assert all(e.metadata.get("parser") == "xls" for e in els)
    sheets = {e.metadata.get("sheet") for e in els}
    assert sheets == {"Sheet1", "Sheet2"}
    header = els[0]
    assert header.text == "Name | Count | Score | Joined | Active"
    assert header.location_label == "XLS Sheet1 row 1"
    data_row = els[1]
    assert data_row.location_label == "XLS Sheet1 row 2"
    assert "123" in data_row.text.split(" | ")   # 整数 cell 无 ".0" 污染
    assert "123.0" not in data_row.text
    assert "1.5" in data_row.text
    assert "2024-03-15T00:00:00" in data_row.text  # 日期 cell ISO 格式
    assert "TRUE" in data_row.text                 # 布尔 cell
    # 空行（Sheet1 第 3 行）被跳过：row_index 从 1（表头）、2（Alice）直接跳到 4（Bob）。
    sheet1_row_indexes = [
        e.metadata["row_index"] for e in els if e.metadata.get("sheet") == "Sheet1"
    ]
    assert sheet1_row_indexes == [1, 2, 4]
    sheet2_texts = [e.text for e in els if e.metadata.get("sheet") == "Sheet2"]
    assert sheet2_texts == ["Notes", "Second sheet row one.", "Second sheet row two."]


def test_parse_source_file_dispatches_xls_without_mineru_client():
    """.xls 不在 MINERU_CAPABLE_SUFFIXES 里，parse_source_file 必须直接走 xlrd，
    绝不构造/调用 mineru_client（旧二进制格式 MinerU 根本不支持）。"""
    client = FakeMineru(content_list=[{"type": "text", "text": "should never be used.", "page_idx": 0}])
    els = parse_source_file("s1", str(_LEGACY_XLS_FIXTURE), "sample_legacy.xls", client)
    assert client.calls == []
    assert all(e.metadata.get("parser") == "xls" for e in els)
    assert len(els) > 0


def test_parse_xls_detects_a_renamed_ooxml_workbook(tmp_path):
    """老系统常把已升级的 .xlsx 改名 .xls 重新分发。这种文件是 ZIP(OOXML) 容器，
    xlrd 会直接抛 XLRDError；ZIP 魔数探测必须把它交给 openpyxl 而不是让整次摄取
    失败。"""
    from openpyxl import Workbook

    real_xlsx = tmp_path / "renamed.xlsx"
    book = Workbook()
    sheet = book.active
    sheet.title = "Sheet1"
    sheet.append(["Header A", "Header B"])
    sheet.append(["value one", "42"])
    book.save(str(real_xlsx))

    fake_xls = tmp_path / "renamed.xls"
    fake_xls.write_bytes(real_xlsx.read_bytes())

    els = parse_xls("s1", fake_xls)
    assert els
    assert all(e.metadata.get("parser") == "xlsx" for e in els)   # openpyxl 形状
    assert any("value one" in e.text for e in els)


# -- _xls_cell_text：纯函数，表驱动钉住 xlrd 归一化契约的每个分支 -------------
#
# stub cell 只需 ctype/value 两个属性——xlrd.Cell 的真实构造依赖打开的 workbook，
# 没必要为了测一个纯函数造一整本工作簿。


class _StubXlsCell:
    def __init__(self, ctype: int, value: Any) -> None:
        self.ctype = ctype
        self.value = value


def test_xls_cell_text_normalizes_every_ctype():
    import xlrd

    from app.services.parsers import _xls_cell_text

    cases = [
        # (ctype, value, datemode, expected)
        (xlrd.XL_CELL_TEXT, "  hello\nworld  ", 0, "hello world"),   # 前后空白 + 内嵌换行压平
        (xlrd.XL_CELL_NUMBER, 123.0, 0, "123"),                       # 整数值float→无".0"
        (xlrd.XL_CELL_NUMBER, 1.5, 0, "1.5"),                          # 非整数浮点原样
        (xlrd.XL_CELL_NUMBER, float("inf"), 0, "inf"),                 # 合法 RK 记录可携带 inf
        (xlrd.XL_CELL_NUMBER, float("nan"), 0, "nan"),                 # NaN 落 str(value)
        (xlrd.XL_CELL_BOOLEAN, 0, 0, "FALSE"),
        (xlrd.XL_CELL_BOOLEAN, 1, 0, "TRUE"),
        (xlrd.XL_CELL_ERROR, 0x2A, 0, "#N/A"),                          # 错误码保留标准错误文本（codex R5 P2）
        (xlrd.XL_CELL_ERROR, 0x07, 0, "#DIV/0!"),                       # 丢成空串会让该格从行拼接里消失、列错位
        (xlrd.XL_CELL_ERROR, 0xEE, 0, "#ERR!"),                         # 未知错误码兜底占位，不再整格蒸发
        (xlrd.XL_CELL_BLANK, "", 0, ""),
        (xlrd.XL_CELL_EMPTY, "", 0, ""),
        # datemode 真的被使用：同一个序列值在 1900/1904 纪元下解出不同日期。
        (xlrd.XL_CELL_DATE, 45366.0, 0, "2024-03-15T00:00:00"),
        (xlrd.XL_CELL_DATE, 45366.0, 1, "2028-03-16T00:00:00"),
    ]
    for ctype, value, datemode, expected in cases:
        rendered = _xls_cell_text(_StubXlsCell(ctype, value), datemode)
        assert rendered == expected, (ctype, value, datemode, rendered)


def test_xls_upload_is_accepted_doc_and_ppt_still_rejected():
    """上传后缀白名单：.xls 放行，.doc/.ppt（无纯 Python 读取器、MinerU 也不支持）
    仍拒绝并提示另存为。"""
    _validate_source_file("sample_legacy.xls")  # 不应抛出
    for rejected_name in ("legacy.doc", "legacy.ppt"):
        with pytest.raises(Exception) as excinfo:
            _validate_source_file(rejected_name)
        assert "Unsupported source file type" in str(getattr(excinfo.value, "detail", excinfo.value))


def test_mineru_label_prefix_maps_upload_suffixes(tmp_path):
    """标签前缀的单一映射点：cloud-only 部署的云端分支与本地 parse_* 分支共用它，
    否则上传 .xlsx 会拿到 "PDF p.1 ..." 与 source_format="pdf"。"""
    assert mineru_label_prefix("book.xlsx") == "XLSX"
    assert mineru_label_prefix("macro.XLSM") == "XLSX"     # 大小写不敏感
    assert mineru_label_prefix("doc.docx") == "DOCX"
    assert mineru_label_prefix("deck.pptx") == "PPTX"
    assert mineru_label_prefix("paper.pdf") == "PDF"
    assert mineru_label_prefix("") == "PDF"                # 缺省/无名 → 历史行为


def test_fallback_warning_suffixes_are_a_strict_subset_of_capable(tmp_path):
    """警告集只覆盖**有损**兜底。xlsx/xlsm 的 openpyxl 兜底对单元格值全保真，
    打「质量降级」是噪音，而且 MinerU 不支持工作簿的部署会得到一个用户永远消不掉
    的假警告。"""
    assert set(MINERU_FALLBACK_WARNING_SUFFIXES) < set(MINERU_CAPABLE_SUFFIXES)
    assert set(MINERU_CAPABLE_SUFFIXES) - set(MINERU_FALLBACK_WARNING_SUFFIXES) == {
        ".xlsx",
        ".xlsm",
    }


# 后缀 → 最小可解析临时文件内容。MinerU 优先分支的后缀不需要真文件（假 client
# 直接返回可用 content_list，本地库解析器根本不会被调到；xlsx 的行覆盖对账读不了
# 桩文件时 fail-open 采信）；没有 MinerU 分支的后缀会真的走本地解析器，所以给的是
# 该格式的合法最小内容。.xls 没有通用桩字节可用（xlrd 需要真正的 BIFF/OLE2
# 结构，通用 zip 桩会直接抛异常），复用 parse_xls 自己的测试 fixture。
_MINIMAL_SOURCE_BYTES = {
    ".md": b"# Title\n\nBody paragraph.\n",
    ".markdown": b"# Title\n\nBody paragraph.\n",
    ".csv": b"a,b\n1,2\n",
    ".xls": _LEGACY_XLS_FIXTURE.read_bytes(),
}


@pytest.mark.parametrize("suffix", sorted(SUPPORTED_SOURCE_SUFFIXES))
def test_mineru_client_reaches_exactly_the_capable_suffixes(tmp_path, suffix):
    """行为守卫：`parse_source_file` 对且仅对 MINERU_CAPABLE_SUFFIXES 里的后缀调用
    MinerU client。断言常量等于某个字面量集合守不住任何东西——真正要拦的是「有人给
    .csv 加了 MinerU 优先分支却没改常量」（云端上传闸会跟着漏发）或反过来「给常量
    加了后缀却没有对应分支」（云端白发一次往返）。"""
    path = tmp_path / f"probe{suffix}"
    path.write_bytes(_MINIMAL_SOURCE_BYTES.get(suffix, b"PK\x03\x04stub"))
    client = FakeMineru(content_list=[{"type": "text", "text": "From MinerU.", "page_idx": 0}])

    _parse_via_chain("s1", path, path.name, client)

    assert bool(client.calls) == (suffix in MINERU_CAPABLE_SUFFIXES)
