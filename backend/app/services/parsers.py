from __future__ import annotations

import base64
import binascii
import html as html_module
import math
import os
import re
import zipfile
from html.parser import HTMLParser
from pathlib import Path, PurePosixPath
from typing import Any, Dict, List
from urllib.parse import unquote
from xml.etree import ElementTree

from app.models.sources import SourceElement
from app.services.knowhow.assets import ALLOWED_MIME_EXTENSIONS
from app.services.parser_registry import PARSER_ENGINES, builtin_parser_id


# A Markdown bundle is one source, so its complete uncompressed file set must
# fit the same deployment byte rail as any other source. The entry ceiling is
# a structural ZIP parser bound rather than a deployment quality/cost knob.
MARKDOWN_BUNDLE_MAX_ENTRIES = 2000
_MARKDOWN_BUNDLE_EXTENSIONS = {".md", ".markdown"}
_MARKDOWN_BUNDLE_COMPRESSION = {zipfile.ZIP_STORED, zipfile.ZIP_DEFLATED}
_REMOTE_OR_DATA_SCHEME = re.compile(r"^[a-z][a-z0-9+.-]*:", re.IGNORECASE)


#: 「MinerU 值得一试」的后缀单一真源。MinerU 原生解析 PDF/DOCX/PPTX/XLSX，
#: 下面这些后缀因此各有一条 MinerU 优先分支（失败/空产出再 fall through 到本地
#: 库兜底）。摄取层用它决定：云端 mineru.net 上传只对这些后缀发起——.md/.csv
#: 这类后缀上传等于把用户内容外发给一个根本解析不了它的第三方，还白付一次重试
#: 预算。
#: .xlsm 与 .xlsx 同为 OOXML 工作簿；远端若不支持会抛异常，按既有降级合同
#: fail-open 回落 openpyxl。
#: ⚠ 这**不是**降级质量警告的判据，那个用下面更窄的
#: MINERU_FALLBACK_WARNING_SUFFIXES。
MINERU_CAPABLE_SUFFIXES = tuple(
    sorted(
        {
            f".{ext}"
            for engine in PARSER_ENGINES
            if engine.id.startswith("mineru_")
            for ext in engine.file_extensions
        }
    )
)

#: 「降级到本机库解析属于有损、值得给用户提示」的后缀。刻意是
#: MINERU_CAPABLE_SUFFIXES 的**真子集**：xlsx/xlsm 的兜底 openpyxl 对单元格值是
#: 全保真的（拿到的是工作簿里逐格的原值，不是版面重建），降级不丢内容，打「质量
#: 降级」警告纯属噪音；而部署的 MinerU 若根本不支持工作簿，那条警告会在每份表格
#: 上永远点亮、重新解析也消不掉——一个用户无从消除的假异常比没有提示更糟。
#: PDF/DOCX/PPTX 的兜底则确实有损（版面、公式、表格、OCR），警告是真信息。
MINERU_FALLBACK_WARNING_SUFFIXES = (".pdf", ".docx", ".pptx")

#: MinerU 工作簿产出的**接受判据**（协议边界，不是可调预算），**行与格两个维度共用
#: 的覆盖阈值**：MinerU 按渲染页出表，受打印区域/列宽/隐藏列影响可能整列整 sheet
#: 丢内容，而且丢了不报错——采信一份残缺产出就是静默数据损失。
#: 只数行是不够的：宽表场景 MinerU 常常**保住每个 `<tr>` 却丢掉页宽之外的列**，
#: 行分子照样足额，静默丢掉大半单元格。因此对账同时算行覆盖与格覆盖，**两个维度
#: 都要达标**才采信，任一低于这个比例就整份丢弃、回落 openpyxl。对账是纯本地流式
#: 计数（openpyxl read_only），零模型零网络；openpyxl 兜底对单元格值全保真，所以
#: 拒收 MinerU 不损失任何数据，只损失表格 HTML 结构——宁可要全量的平铺行，不要残
#: 缺的结构化表。
MINERU_WORKBOOK_MIN_ROW_COVERAGE = 0.8

#: 需要跑上面这道对账的后缀。摄取层的云端**文件**分支与本地 parse_xlsx 共用它，
#: 避免 cloud-only 部署绕过对账直接采信残缺产出。
MINERU_WORKBOOK_SUFFIXES = (".xlsx", ".xlsm")

#: 本机库兜底路径（mammoth / python-pptx）单个 table 元素最多容纳的物理行数。
#: 这是**检索粒度边界**不是显示预算：`chunking.py` 从不切分单个元素，所以一张 1500
#: 行的 Word 表折成一个元素就是一个 20 万字符的 chunk——embedding 截断 + 召回坍缩，
#: 同一份内容旧 python-docx 兜底能出 300 个可检索单位。MinerU 路径不适用：它按渲染
#: 页出表，每页一张表天然有界。
_FALLBACK_TABLE_ROWS_PER_ELEMENT = 30

#: MinerU 元素 location_label / metadata.source_format 的前缀。按上传文件后缀映射，
#: 未知后缀落 "PDF"（MinerU 的默认与历史行为）。本地 parse_* 分支与摄取层的云端
#: 分支共用它，避免 cloud-only 部署上传 .xlsx 却拿到 "PDF p.1 ..." 这种错标。
_MINERU_LABEL_PREFIXES = {
    ".docx": "DOCX",
    ".pptx": "PPTX",
    ".xlsx": "XLSX",
    ".xlsm": "XLSX",
}


def mineru_label_prefix(file_name: str) -> str:
    """按文件名后缀给出 MinerU 元素标签前缀（未知后缀→ "PDF"）。"""
    return _MINERU_LABEL_PREFIXES.get(Path(file_name or "").suffix.lower(), "PDF")


def parse_builtin_source_file(
    source_id: str,
    file_path: str,
    file_name: str,
    persist_image: Any = None,
) -> List[SourceElement]:
    """Dispatch only to the guaranteed local parser surface.

    Optional MinerU selection belongs exclusively to the parser ProviderChain.
    Keeping this function provider-free prevents a builtin fallback from
    re-entering an earlier link and issuing a second remote request.
    """
    parser_id = builtin_parser_id(file_name)
    if parser_id == "markdown":
        return parse_markdown(source_id, Path(file_path), persist_image=persist_image)
    if parser_id == "markdown_bundle":
        return parse_markdown_bundle(
            source_id, Path(file_path), persist_image=persist_image
        )
    if parser_id == "docx":
        return parse_docx_mammoth(source_id, Path(file_path))
    if parser_id == "pptx":
        return parse_pptx_pythonpptx(source_id, Path(file_path))
    if parser_id == "pdf":
        return parse_pdf_python(source_id, Path(file_path))
    if parser_id == "csv":
        return parse_csv(source_id, Path(file_path))
    if parser_id == "xlsx":
        return parse_xlsx_basic(source_id, Path(file_path))
    if parser_id == "xls":
        # 旧版二进制 Excel（BIFF）。MinerU 不支持该格式，没有优先分支——不进
        # MINERU_CAPABLE_SUFFIXES，也不传 mineru_client。xlrd 是纯 Python 读取器，
        # 直接兜底到底。
        return parse_xls(source_id, Path(file_path))
    return parse_plain_text(source_id, Path(file_path), "text")


def parse_source_file(
    source_id: str,
    file_path: str,
    file_name: str,
    mineru_client: Any = None,
    persist_image: Any = None,
) -> List[SourceElement]:
    """Backward-compatible name for the provider-free builtin dispatcher."""

    return parse_builtin_source_file(
        source_id, file_path, file_name, persist_image=persist_image
    )


def parse_csv(source_id: str, path: Path) -> List[SourceElement]:
    import csv

    elements: List[SourceElement] = []
    with path.open(newline="", encoding="utf-8", errors="replace") as handle:
        for index, row in enumerate(csv.reader(handle), start=1):
            cells = [cell.strip() for cell in row if cell and cell.strip()]
            if not cells:
                continue
            elements.append(
                _element(
                    source_id,
                    "table_row",
                    f"CSV row {index}",
                    " | ".join(cells),
                    {"parser": "csv", "row_index": index},
                )
            )
    return elements


def parse_xlsx(
    source_id: str,
    path: Path,
    file_name: str = "",
    mineru_client: Any = None,
    persist_image: Any = None,
) -> List[SourceElement]:
    """Compatibility wrapper for the provider-free openpyxl fallback."""
    return parse_xlsx_basic(source_id, path)


def parse_xlsx_basic(source_id: str, path: Path) -> List[SourceElement]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("XLSX parser dependency openpyxl is not installed") from exc

    elements: List[SourceElement] = []
    # Open via a binary file handle rather than `str(path)`: openpyxl's
    # `_validate_archive` gates filename inputs on file extension (and
    # explicitly rejects ".xls" with a dedicated error message) but skips that
    # check entirely for file-like objects. parse_xls() routes ZIP-magic
    # ".xls" files (renamed .xlsx redistributions) here, so a bare filename
    # string would always raise InvalidFileException regardless of the actual
    # ZIP/OOXML content.
    with path.open("rb") as handle:
        workbook = load_workbook(handle, read_only=True, data_only=True)
        for sheet in workbook.worksheets:
            for index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if not cells:
                    continue
                elements.append(
                    _element(
                        source_id,
                        "table_row",
                        f"XLSX {sheet.title} row {index}",
                        " | ".join(cells),
                        {"parser": "xlsx", "sheet": sheet.title, "row_index": index},
                    )
                )
        workbook.close()
    return elements


def parse_xls(source_id: str, path: Path) -> List[SourceElement]:
    """Parse a legacy binary .xls (BIFF) workbook via xlrd.

    MinerU cannot parse the old OLE2 binary format at all, so this has no
    MinerU-first branch — xlrd (pure Python) is the only reader, mirroring the
    row-dump shape of parse_xlsx_basic (openpyxl) for output consistency, with
    the closest xlrd offers to openpyxl's read_only streaming: `on_demand=True`
    loads one sheet at a time and each sheet is unloaded via
    `book.unload_sheet()` right after it is consumed (xlrd has no true
    row-level streaming, so per-sheet load/unload is the best it can do).
    Numeric cells come back as float from xlrd; integer-valued floats are
    rendered without a trailing ".0" so retrieval text isn't polluted with
    "123.0". Date cells are converted via the workbook's datemode. Error cells
    are skipped (no usable text).

    Some deployments rename an already-migrated .xlsx to .xls when
    redistributing older assets. Such a file is a ZIP (OOXML) container, not
    BIFF/OLE2, and xlrd raises XLRDError on it. Detect the ZIP magic up front
    and hand it to the OOXML parser instead of failing outright.
    """
    if zipfile.is_zipfile(path):
        return parse_xlsx_basic(source_id, path)

    try:
        import xlrd
    except ImportError as exc:
        raise RuntimeError("XLS parser dependency xlrd is not installed") from exc

    elements: List[SourceElement] = []
    # xlrd's per-record diagnostic logging (`open_workbook`'s default
    # `logfile=sys.stdout`) is not gated by any verbosity flag and can echo
    # fragments of the workbook's raw bytes into stdout — which lands in
    # backend.log unbounded for a corrupt/trailing-garbage .xls. Discard it
    # rather than collect it: collecting into memory just swaps a log flood
    # for a memory flood.
    with open(os.devnull, "w") as devnull:
        book = xlrd.open_workbook(str(path), on_demand=True, logfile=devnull)
        try:
            for name in book.sheet_names():
                sheet = book.sheet_by_name(name)
                for index in range(sheet.nrows):
                    cells: List[str] = []
                    for cell in sheet.row(index):
                        rendered = _xls_cell_text(cell, book.datemode)
                        if rendered:
                            cells.append(rendered)
                    if cells:
                        elements.append(
                            _element(
                                source_id,
                                "table_row",
                                f"XLS {name} row {index + 1}",
                                " | ".join(cells),
                                {"parser": "xls", "sheet": name, "row_index": index + 1},
                            )
                        )
                book.unload_sheet(name)
        finally:
            book.release_resources()
    return elements


def _xls_cell_text(cell: Any, datemode: int) -> str:
    """Render one xlrd Cell to display text, or "" if it carries nothing usable."""
    import xlrd

    if cell.ctype == xlrd.XL_CELL_EMPTY or cell.ctype == xlrd.XL_CELL_BLANK:
        return ""
    if cell.ctype == xlrd.XL_CELL_ERROR:
        # `#N/A` / `#DIV/0!` 是工作簿的真实内容（codex R5 P2）：丢成空串会让该格
        # 从 " | " 拼接里整个消失（列错位），全错误行更是整行蒸发。openpyxl 路径
        # (data_only) 保留的正是这些标准错误文本，两条工作簿路径口径要一致。
        try:
            from xlrd.biffh import error_text_from_code

            text = error_text_from_code.get(int(cell.value), "")
        except Exception:
            text = ""
        return text or "#ERR!"
    if cell.ctype == xlrd.XL_CELL_BOOLEAN:
        return "TRUE" if cell.value else "FALSE"
    if cell.ctype == xlrd.XL_CELL_DATE:
        try:
            return xlrd.xldate_as_datetime(cell.value, datemode).isoformat()
        except Exception:
            return str(cell.value).strip()
    if cell.ctype == xlrd.XL_CELL_NUMBER:
        value = cell.value
        # A legitimate RK-encoded numeric cell can carry inf/NaN; `int(value)`
        # on either raises OverflowError/ValueError, and letting that escape
        # would blow up ingestion of the whole workbook over a single cell.
        if math.isfinite(value) and value == int(value):
            return str(int(value))
        return str(value).strip()
    # Text cells (and any other/unknown ctype): normalize embedded whitespace
    # too, not just leading/trailing — a cell with an Alt+Enter line break
    # would otherwise leak a raw "\n" into the joined " | "-delimited row.
    return " ".join(str(cell.value).split())


def _workbook_nonempty_counts(path: Path) -> tuple[int, int] | None:
    """工作簿全部 sheet 的 (非空行总数, 非空格总数)，MinerU 产出对账的两个分母。

    「非空行」判据与 parse_xlsx_basic 的过滤逐字一致（至少一个非 None 且 strip
    非空的 cell），「非空格」是同一判据逐格版本，否则分子分母不同口径、对账无意
    义。两个分母在**同一次**流式遍历里数出来（read_only，不物化整个工作簿，也不
    为第二个指标再扫一遍）。openpyxl 缺失或读失败返回 None——对账绝不引入新的失
    败面，调用方按 fail-open 采信 MinerU。
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        return None
    workbook = None
    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
        total_rows = 0
        total_cells = 0
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                filled = sum(
                    1 for cell in row if cell is not None and str(cell).strip()
                )
                if filled:
                    total_rows += 1
                    total_cells += filled
        return total_rows, total_cells
    except Exception:
        return None
    finally:
        if workbook is not None:
            try:
                workbook.close()
            except Exception:
                pass


_TABLE_ROW_SEGMENT_RE = re.compile(r"<tr\b.*?(?:</tr>|(?=<tr\b)|$)", re.IGNORECASE | re.DOTALL)
_TABLE_CELL_BODY_RE = re.compile(r"<t[dh]\b[^>]*>(.*?)</t[dh]>", re.IGNORECASE | re.DOTALL)


def _mineru_workbook_coverage(elements: List[SourceElement]) -> tuple[int, int]:
    """MinerU 产出覆盖了多少工作簿 (非空行, 非空格)——对账的两个分子。

    分子与分母（openpyxl 的非空判定）必须是**同一套非空语义**（codex R6 P1）：
    只按 `<tr`/`<td` **标签**计数的话，MinerU 返回一张结构完整、值却全空的格子
    （识别失败时的常见形态）也能刷满两个阈值，把全保真的 openpyxl 兜底绕过去——
    护栏本身被掏空。因此：

    - 表格元素逐 `<tr>` 段扫描：一个格只有在剥标签 + 反转义后仍有实文时才计入
      格分子；一行只有在至少含一个这样的非空格时才计入行分子。空 `<td></td>`、
      空 `<tr></tr>` 一律不算覆盖（大小写不敏感——MinerU 也吐过 `<TR>`）。
    - 非表格**文本**元素每条只向两个分子**各计 1**：MinerU 会把 sheet 标题、游离
      单元格、图注映射成标题/段落而不是表格行，不计入会把正常产出误判成缺行；但
      一条段落最多只能顶一行一格，虚增因此有界（元素映射层已丢弃空文本元素）。
    - image 元素**不计入任何分子**：图片不对应任何物理行或格，计入它等于让一份
      「几张图 + 少量行」的残缺产出把分子刷够。
    """
    rows = 0
    cells = 0
    for element in elements:
        if element.element_type == "table":
            html = str(element.metadata.get("table_html", "") or "")
            for row_segment in _TABLE_ROW_SEGMENT_RE.findall(html):
                populated = 0
                for cell_body in _TABLE_CELL_BODY_RE.findall(row_segment):
                    text = html_module.unescape(re.sub(r"<[^>]+>", " ", cell_body))
                    if text.strip():
                        populated += 1
                cells += populated
                if populated:
                    rows += 1
        elif element.element_type == "image":
            continue
        else:
            rows += 1
            cells += 1
    return rows, cells


def _mineru_workbook_reconcile(
    path: Path, elements: List[SourceElement]
) -> tuple[bool, int, int, int, int]:
    """(是否采信, 覆盖行数, 非空行总数, 覆盖格数, 非空格总数)。

    见 MINERU_WORKBOOK_MIN_ROW_COVERAGE：行覆盖与格覆盖**都**达标才采信。分母为 0
    （空工作簿）或对账本身不可用时一律 fail-open 采信——对账是护栏，不能自己变成
    新的失败模式。
    """
    try:
        counts = _workbook_nonempty_counts(path)
        if counts is None:
            return True, 0, 0, 0, 0
        total_rows, total_cells = counts
        if not total_rows:
            return True, 0, 0, 0, 0
        rows, cells = _mineru_workbook_coverage(elements)
        accepted = rows >= math.ceil(
            total_rows * MINERU_WORKBOOK_MIN_ROW_COVERAGE
        ) and cells >= math.ceil(total_cells * MINERU_WORKBOOK_MIN_ROW_COVERAGE)
        return accepted, rows, total_rows, cells, total_cells
    except Exception:
        return True, 0, 0, 0, 0


def mineru_workbook_output_accepted(path: Path, elements: List[SourceElement]) -> bool:
    """工作簿产出对账的公共入口（摄取层的云端文件分支复用）。

    云端 cloud-only 部署走的是 source_ingestion 自己的映射路径，不经 parse_xlsx；
    没有这个入口它就会绕过对账直接采信残缺产出。
    """
    return _mineru_workbook_reconcile(path, elements)[0]


def mineru_workbook_reconciliation(
    path: Path, elements: List[SourceElement]
) -> tuple[bool, int, int, int, int]:
    """Expose the pure reconciliation receipt to the core parser adapter."""

    return _mineru_workbook_reconcile(path, elements)


def _persist_markdown_data_uri(src: str, persist_image: Any, ordinal: int) -> str:
    """解析 `data:<mime>;base64,<payload>` 并按 MinerU 同款契约持久化。

    任何一步不合规（无 persist_image、非 base64 编码、mime 不在白名单、解码
    失败、persist 本身失败）都返回空串——图片问题绝不阻塞文本解析。单图字节
    上限/每源张数上限由 persist_image 闭包内置护栏负责，这里不新增数字上限。
    """
    if persist_image is None:
        return ""
    if "," not in src:
        return ""
    header, _, payload = src.partition(",")
    if not payload:
        return ""
    header_body = header[len("data:"):]
    parts = header_body.split(";")
    mime = (parts[0] or "").strip().lower()
    if "base64" not in [p.strip().lower() for p in parts[1:]]:
        return ""
    if mime not in ALLOWED_MIME_EXTENSIONS:
        return ""
    # Decode happens before persist_image's own byte-size guard runs, so a
    # single oversized data URI is fully decoded in memory first — accepted
    # here rather than pre-checking encoded length: the upload path already
    # has SOURCE_UPLOAD_MAX_MB as an outer bound, the offline batch_ingest
    # path is operator-controlled input, and the allocation is transient
    # (freed once persist_image returns or this function returns "").
    try:
        img_bytes = base64.b64decode(payload, validate=True)
    except (binascii.Error, ValueError):
        return ""
    ext = ALLOWED_MIME_EXTENSIONS[mime]
    try:
        asset_id = persist_image(img_bytes, f"md-img-{ordinal}.{ext}")
    except Exception:
        return ""  # 抽图/写盘失败绝不影响文本产出
    return asset_id or ""


def parse_markdown_text(
    source_id: str,
    text: str,
    persist_image: Any = None,
    resolve_image: Any = None,
) -> List[SourceElement]:
    """parse_markdown 的无文件版本：直接解析给定 markdown 文本（Memory 派生源等
    没有磁盘文件的调用方复用）。逐字复用 parse_markdown 原先内嵌的 parse_blocks
    调用与元素装配逻辑，仅将「从文件读文本」换成「拿到手的文本」。

    `persist_image(img_bytes, img_name) -> Optional[str]` 镜像 MinerU 路径的
    persist 闭包，用于把 data URI 内嵌的图片落盘为资产；省略时 data URI 被剥
    离而不入库（图片不可用，但文本解析不受影响）。
    """
    from app.services.structural_markdown import parse_blocks

    blocks = parse_blocks(text)
    elements: List[SourceElement] = []
    counters: Dict[str, int] = {}
    for block in blocks:
        if block.type == "image_description":
            # `> **图片描述**` 引用块的文本已经折进上一条 image 元素（见下面的
            # image 分支），再出一条元素就是把同一段描述发两遍——两条都会进
            # chunk，同一份内容占两个检索位。它仍以独立块的形式存在，是为了让
            # KG 侧拿到逐字原文跨度。
            continue
        counters[block.type] = counters.get(block.type, 0) + 1
        ordinal = counters[block.type]
        metadata: Dict[str, Any] = {
            "parser": "markdown",
            "section_path": block.section_path,
            "char_start": block.char_start,
            "char_end": block.char_end,
            "line_start": block.line_start,
            "line_end": block.line_end,
        }
        if block.type == "heading":
            metadata["heading_level"] = block.level
            if block.anchor_id:
                metadata["anchor_id"] = block.anchor_id
        if block.type == "code_block":
            metadata["lang"] = block.lang
        if block.type == "table" and block.raw.lstrip().lower().startswith("<table"):
            # PyMuPDF4LLM can emit reconstructed HTML tables in reading order.
            # Preserve that renderable representation while Block.text remains
            # the compact retrieval/embedding form.
            metadata["table_html"] = block.raw.strip()
        element_text = block.text
        if block.type == "image":
            caption = block.text.strip()
            if caption:
                metadata["caption"] = caption
            # 紧跟这张图的 `> **图片描述**` 引用块（structural_markdown 已折进
            # block.metadata）。它与 alt 图注并列、不互相覆盖：alt 是图注，描述
            # 是对图的展开说明，界面分两处渲染，检索则两段都要。
            raw_description = block.metadata.get("description", "")
            description = raw_description.strip() if isinstance(raw_description, str) else ""
            if description:
                metadata["description"] = description
            src = block.metadata.get("src", "")
            asset_id = ""
            # Scheme match is case-insensitive: markdown-it tokenizes
            # `DATA:image/png;base64,...` as an image, and a case-sensitive
            # check here would dump the whole base64 URI into metadata["src"].
            if isinstance(src, str) and src[:5].lower() == "data:":
                # data URI 绝不进 metadata（任何键都不行）：只尝试持久化。
                asset_id = _persist_markdown_data_uri(src, persist_image, ordinal)
                if asset_id:
                    metadata["asset_id"] = asset_id
            elif isinstance(src, str) and resolve_image is not None:
                try:
                    asset_id = resolve_image(src, ordinal) or ""
                except Exception:
                    asset_id = ""
                if asset_id:
                    metadata["asset_id"] = asset_id
                else:
                    metadata["src"] = src
            else:
                metadata["src"] = src
            if not caption and not description and not asset_id:
                continue  # 既不可显示又不可检索的占位行没有价值
            # 图注与描述都进元素文本：元素文本就是这张图进 chunk 的那段字，少
            # 带一半等于让描述检索不到。分隔符是换行只为可读——`_element` 对
            # image 类型照旧压平空白，落库的是单行。
            element_text = "\n".join(
                part for part in (caption, description) if part
            ) or f"Markdown 图 {ordinal}"
        elements.append(
            _element(
                source_id,
                block.type,
                f"Markdown {block.type} {ordinal}",
                element_text,
                metadata,
            )
        )
    if blocks:
        # A discarded (uncaptioned, unpersisted) data-URI image can leave
        # `elements` empty while `blocks` is not; return the (possibly empty)
        # element list rather than falling through to the raw-text fallback,
        # which would otherwise dump the giant base64 payload into an element.
        return elements
    # No structured blocks (e.g. blank input): same blank-line paragraph
    # fallback as parse_plain_text, minus the file read. The raw text is
    # sanitized first: a document whose only content is alt-less rejected-mime
    # data-URI image literals inside containers `parse_blocks` drops entirely
    # (e.g. a lone `- ![](data:image/svg+xml;base64,...)` list item) reaches
    # this branch, and splitting the unsanitized original would dump the
    # base64 payload into a paragraph element.
    from app.services.structural_markdown import strip_data_uri_image_literals

    text = strip_data_uri_image_literals(text)
    chunks = [chunk.strip() for chunk in re.split(r"\n\s*\n", text) if chunk.strip()]
    return [
        _element(
            source_id,
            "paragraph",
            f"Text paragraph {index}",
            " ".join(chunk.split()),
            {"parser": "markdown", "paragraph_index": index},
        )
        for index, chunk in enumerate(chunks, start=1)
    ]


def parse_markdown(
    source_id: str, path: Path, persist_image: Any = None
) -> List[SourceElement]:
    return parse_markdown_text(
        source_id,
        path.read_text(encoding="utf-8", errors="replace"),
        persist_image=persist_image,
    )


def _normalize_bundle_entry_path(raw: str) -> str | None:
    """Return a safe, root-relative POSIX archive path or ``None``.

    ZIP member names are attacker-controlled. Backslashes are separators for
    exported Windows bundles, while absolute paths, NULs and ``..`` escaping
    the archive root are rejected before any member is read.
    """
    if not isinstance(raw, str) or not raw or "\x00" in raw:
        return None
    replaced = raw.replace("\\", "/")
    if replaced.startswith("/"):
        return None
    parts: list[str] = []
    for part in replaced.split("/"):
        if part in {"", "."}:
            continue
        if part == "..":
            if not parts:
                return None
            parts.pop()
            continue
        parts.append(part)
    return "/".join(parts) or None


def _bundle_image_path(markdown_path: str, src: str) -> str | None:
    """Resolve one Markdown image target against its owning bundle document."""
    value = src.strip()
    if not value or value.startswith("//") or _REMOTE_OR_DATA_SCHEME.match(value):
        return None
    # Query/fragment suffixes are URL syntax, not part of the stored file name.
    value = value.split("#", 1)[0].split("?", 1)[0]
    try:
        value = unquote(value)
    except (UnicodeDecodeError, ValueError):
        return None
    if value.startswith("/"):
        combined = value.lstrip("/")
    else:
        parent = str(PurePosixPath(markdown_path).parent)
        combined = value if parent == "." else f"{parent}/{value}"
    return _normalize_bundle_entry_path(combined)


def _bundle_image_extension(data: bytes) -> str:
    """Sniff the four source-asset raster formats from magic bytes."""
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return "png"
    if data.startswith(b"\xff\xd8\xff"):
        return "jpg"
    if data.startswith((b"GIF87a", b"GIF89a")):
        return "gif"
    if len(data) >= 12 and data.startswith(b"RIFF") and data[8:12] == b"WEBP":
        return "webp"
    return ""


def parse_markdown_bundle(
    source_id: str,
    path: Path,
    persist_image: Any = None,
    *,
    max_uncompressed_bytes: int | None = None,
) -> List[SourceElement]:
    """Parse every Markdown document in a ZIP and persist relative images.

    The raw archive remains the source of record. Parsing is background work,
    and relative references are resolved against each Markdown member's own
    directory. Missing/remote/unsupported images fail open as caption text;
    malformed or unsafe archive structure rejects the source as a whole.
    """
    from app.core.config import get_settings

    byte_limit = (
        get_settings().source_upload_max_bytes
        if max_uncompressed_bytes is None
        else max_uncompressed_bytes
    )
    if byte_limit <= 0:
        raise ValueError("Markdown bundle byte limit must be positive")

    try:
        archive = zipfile.ZipFile(path)
    except (OSError, zipfile.BadZipFile) as exc:
        raise ValueError("invalid Markdown ZIP bundle") from exc

    with archive:
        members: dict[str, zipfile.ZipInfo] = {}
        total_bytes = 0
        for info in archive.infolist():
            normalized = _normalize_bundle_entry_path(info.filename)
            if normalized is None:
                if info.is_dir():
                    continue
                raise ValueError("Markdown ZIP bundle contains an unsafe path")
            if info.is_dir() or normalized.startswith("__MACOSX/"):
                continue
            if info.flag_bits & 0x1:
                raise ValueError("encrypted Markdown ZIP bundles are not supported")
            if info.compress_type not in _MARKDOWN_BUNDLE_COMPRESSION:
                raise ValueError("Markdown ZIP bundle uses unsupported compression")
            if normalized in members:
                raise ValueError("Markdown ZIP bundle contains duplicate paths")
            members[normalized] = info
            if len(members) > MARKDOWN_BUNDLE_MAX_ENTRIES:
                raise ValueError("Markdown ZIP bundle contains too many files")
            total_bytes += info.file_size
            if total_bytes > byte_limit:
                raise ValueError("Markdown ZIP bundle is too large after decompression")

        markdown_paths = sorted(
            name
            for name in members
            if PurePosixPath(name).suffix.lower() in _MARKDOWN_BUNDLE_EXTENSIONS
        )
        if not markdown_paths:
            raise ValueError("Markdown ZIP bundle contains no Markdown document")

        def read_member(info: zipfile.ZipInfo) -> bytes:
            """Read no more than the central-directory declaration plus one.

            The aggregate admission above is useful only if a corrupt local
            header cannot make ``ZipFile.read`` inflate unbounded bytes before
            the size disagreement is noticed.
            """
            try:
                with archive.open(info) as handle:
                    data = handle.read(info.file_size + 1)
                    trailing = handle.read(1)
            except (OSError, RuntimeError, zipfile.BadZipFile) as exc:
                raise ValueError("Markdown ZIP bundle is corrupt") from exc
            if len(data) != info.file_size or trailing:
                raise ValueError("Markdown ZIP bundle member size is invalid")
            return data

        elements: List[SourceElement] = []
        for markdown_path in markdown_paths:
            markdown_text = read_member(members[markdown_path]).decode(
                "utf-8", errors="replace"
            )

            def resolve_image(src: str, ordinal: int) -> str:
                if persist_image is None:
                    return ""
                image_path = _bundle_image_path(markdown_path, src)
                info = members.get(image_path or "")
                if info is None:
                    return ""
                try:
                    data = read_member(info)
                except ValueError:
                    return ""
                extension = _bundle_image_extension(data)
                if not extension:
                    return ""
                return persist_image(
                    data,
                    f"bundle-img-{ordinal}.{extension}",
                ) or ""

            parsed = parse_markdown_text(
                source_id,
                markdown_text,
                persist_image=persist_image,
                resolve_image=resolve_image,
            )
            for element in parsed:
                element.metadata = {
                    **element.metadata,
                    "bundle_path": markdown_path,
                }
                element.location_label = f"{markdown_path} · {element.location_label}"
            elements.extend(parsed)
        return elements


def parse_plain_text(source_id: str, path: Path, parser_name: str) -> List[SourceElement]:
    text = path.read_text(encoding="utf-8", errors="replace")
    chunks = [chunk.strip() for chunk in re.split(r"\n\s*\n", text) if chunk.strip()]
    return [
        _element(
            source_id,
            "paragraph",
            f"Text paragraph {index}",
            " ".join(chunk.split()),
            {"parser": parser_name, "paragraph_index": index},
        )
        for index, chunk in enumerate(chunks, start=1)
    ]


def parse_docx(
    source_id: str,
    path: Path,
    file_name: str = "",
    mineru_client: Any = None,
    persist_image: Any = None,
) -> List[SourceElement]:
    """Compatibility wrapper for the provider-free DOCX fallback."""
    return parse_docx_mammoth(source_id, path)


#: mammoth 输出里构成一个「块」的标签。表格单独处理（它自己就是一个块，内部的
#: <p> 不是块边界）。
_DOCX_HEADING_TAGS = {"h1", "h2", "h3", "h4", "h5", "h6"}
_DOCX_PARAGRAPH_TAGS = {"p", "li", "blockquote", "pre"}


class _MammothHtmlBlocks(HTMLParser):
    """把 mammoth 的语义化 HTML 归并成块级 (kind, level, text, raw_html) 序列。

    刻意只用 stdlib：本仓库不为一次兜底解析引入 beautifulsoup/markdownify。
    表格子树按原样重新序列化进 raw_html（渲染端认 metadata.table_html），表格内的
    <p> 不当块边界。<img> 标签一律丢弃、只保留 alt 文本并入所在段落——图片字节在
    上游就已被 parse_docx_mammoth 的 convert_image 挡掉（见那里的注释），这里是第
    二道闸：即便某条路径换回默认 handler，base64 data URI 也不会进 text/metadata。
    """

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.blocks: List[tuple[str, int, str, str]] = []
        self._buffer: List[str] = []
        self._kind = ""
        self._level = 0
        self._table_depth = 0
        self._table_raw: List[str] = []
        # 每层列表一个计数器：None = 无序列表，int = 有序列表的当前序号。
        self._list_stack: List[int | None] = []

    # -- 块装配 ---------------------------------------------------------
    def _flush(self) -> None:
        text = " ".join("".join(self._buffer).split())
        if self._kind and text:
            self.blocks.append((self._kind, self._level, text, ""))
        self._buffer = []
        self._kind = ""
        self._level = 0

    def _start_block(self, kind: str, level: int = 0, prefix: str = "") -> None:
        self._flush()
        self._kind = kind
        self._level = level
        self._buffer = [prefix] if prefix else []

    def _list_prefix(self) -> str:
        if not self._list_stack:
            return "- "
        counter = self._list_stack[-1]
        if counter is None:
            return "- "
        counter += 1
        self._list_stack[-1] = counter
        return f"{counter}. "

    # -- HTMLParser 钩子 -------------------------------------------------
    def handle_startendtag(self, tag: str, attrs) -> None:
        # 默认实现会再调一次 handle_endtag，对 <img />/<br /> 这类空元素会在表格
        # raw 里留下伪造的 </img>。只走 start 分支。
        self.handle_starttag(tag, attrs)

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag == "img":
            alt = " ".join(
                " ".join(str(value or "").split())
                for key, value in attrs
                if key == "alt"
            ).strip()
            if alt:
                if self._table_depth:
                    self._table_raw.append(html_module.escape(f" {alt} "))
                else:
                    self._buffer.append(f" {alt} ")
            return
        if self._table_depth:
            if tag == "table":
                self._table_depth += 1
            self._table_raw.append(self.get_starttag_text() or f"<{tag}>")
            return
        if tag == "table":
            self._flush()
            self._table_depth = 1
            self._table_raw = [self.get_starttag_text() or "<table>"]
            return
        if tag in ("ul", "ol"):
            self._list_stack.append(0 if tag == "ol" else None)
            return
        if tag == "br":
            self._buffer.append(" ")
            return
        if tag in _DOCX_HEADING_TAGS:
            self._start_block("heading", int(tag[1]))
            return
        if tag == "li":
            self._start_block("paragraph", 0, self._list_prefix())
            return
        if tag in _DOCX_PARAGRAPH_TAGS:
            self._start_block("paragraph")

    def handle_endtag(self, tag: str) -> None:
        if tag in ("img", "br"):
            return
        if self._table_depth:
            self._table_raw.append(f"</{tag}>")
            if tag == "table":
                self._table_depth -= 1
                if self._table_depth == 0:
                    raw = "".join(self._table_raw)
                    self._table_raw = []
                    text = _html_table_to_text(raw)
                    if text:
                        self.blocks.append(("table", 0, text, raw))
            return
        if tag in ("ul", "ol"):
            if self._list_stack:
                self._list_stack.pop()
            return
        if tag in _DOCX_HEADING_TAGS or tag in _DOCX_PARAGRAPH_TAGS:
            self._flush()

    def handle_data(self, data: str) -> None:
        if self._table_depth:
            self._table_raw.append(html_module.escape(data))
        else:
            self._buffer.append(data)

    def close(self) -> None:
        super().close()
        if self._table_depth:
            # 防御性收尾：mammoth 自己产出良构 HTML，但若 </table> 缺失，累积在
            # _table_raw 里的整张表连同其后的正文会被一起吞掉、一个块都不产出。
            raw = "".join(self._table_raw)
            self._table_raw = []
            self._table_depth = 0
            text = _html_table_to_text(raw)
            if text:
                self.blocks.append(("table", 0, text, raw))
        self._flush()


def _mammoth_html_blocks(markup: str) -> List[tuple[str, int, str, str]]:
    parser = _MammothHtmlBlocks()
    parser.feed(markup)
    parser.close()
    return parser.blocks


def parse_docx_mammoth(source_id: str, path: Path) -> List[SourceElement]:
    """Semantic DOCX fallback: mammoth HTML → structured elements.

    python-docx (parse_docx_basic) walks paragraphs as flat `.text`, losing
    heading levels, list markers and table structure. mammoth maps Word styles
    onto semantic HTML, which we fold into the repo's element vocabulary
    (heading/paragraph/table). Its markdown writer is deprecated and has no
    table support, hence the HTML route.

    Deliberately fail-open: a missing wheel, a conversion error or an empty
    result falls through to parse_docx_basic, which is the last resort.
    """
    try:
        import mammoth

        with path.open("rb") as handle:
            # result.messages are style warnings ("unrecognized paragraph
            # style"); they never reach the user and are intentionally ignored.
            #
            # convert_image is not cosmetic: mammoth's default handler
            # base64-materialises every embedded image into the HTML (a 2.4 MB
            # picture referenced six times became a 19 MB string / 68 MB peak
            # RSS in a local run) and this parser discards all of it. Returning
            # an empty attribute dict keeps the <img> alt text and never reads
            # the image bytes.
            result = mammoth.convert_to_html(
                handle, convert_image=mammoth.images.img_element(lambda image: {})
            )
        blocks = _mammoth_html_blocks(str(result.value or ""))
        elements: List[SourceElement] = []
        counters: Dict[str, int] = {}
        for kind, level, text, raw in blocks:
            counters[kind] = counters.get(kind, 0) + 1
            if kind == "table" and raw:
                parts = _split_table_html(raw)
                for part_index, part in enumerate(parts, start=1):
                    part_text = _html_table_to_text(part)
                    if not part_text:
                        continue
                    label = f"DOCX table {counters[kind]}"
                    metadata: Dict[str, Any] = {
                        "parser": "mammoth",
                        "table_index": counters[kind],
                        "table_html": part,
                    }
                    if len(parts) > 1:
                        label = f"{label} part {part_index}"
                        metadata["table_part"] = part_index
                        metadata["table_parts"] = len(parts)
                    elements.append(_element(source_id, kind, label, part_text, metadata))
                continue
            metadata = {"parser": "mammoth", f"{kind}_index": counters[kind]}
            if kind == "heading":
                metadata["heading_level"] = level
            elements.append(
                _element(
                    source_id,
                    kind,
                    f"DOCX {kind} {counters[kind]}",
                    text,
                    metadata,
                )
            )
        if elements:
            return elements
    except Exception:
        # Missing wheel / unsupported document / parser error: one final try
        # with python-docx so ingestion still completes.
        pass
    return parse_docx_basic(source_id, path)


def parse_docx_basic(source_id: str, path: Path) -> List[SourceElement]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("DOCX parser dependency python-docx is not installed") from exc

    document = Document(str(path))
    elements: List[SourceElement] = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = paragraph.text.strip()
        if not text:
            continue
        elements.append(
            _element(
                source_id,
                "paragraph",
                f"DOCX paragraph {index}",
                text,
                {"parser": "docx", "paragraph_index": index},
            )
        )

    table_index = 1
    for table in document.tables:
        for row_index, row in enumerate(table.rows, start=1):
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if not values:
                continue
            elements.append(
                _element(
                    source_id,
                    "table_row",
                    f"DOCX table {table_index} row {row_index}",
                    " | ".join(values),
                    {"parser": "docx", "table_index": table_index, "row_index": row_index},
                )
            )
        table_index += 1
    return elements


def parse_pptx(
    source_id: str,
    path: Path,
    file_name: str = "",
    mineru_client: Any = None,
    persist_image: Any = None,
) -> List[SourceElement]:
    """Compatibility wrapper for the provider-free PPTX fallback."""
    return parse_pptx_pythonpptx(source_id, path)


#: 组合形状递归的深度上界。纯护栏：PPTX 的 spTree 是有限 XML 树，正常演示文稿最多
#: 嵌两三层；这个数只用来保证「文件损坏/被构造出病态深度」时不会栈溢出，绝不是产品
#: 层面的取舍。旧实现只展开一层，两层分组里的文本会整块丢失，还不如它要取代的原始
#: XML 抽取（`.//p:sp` 是任意深度）。
def _pptx_iter_shapes(shapes: Any):
    """按文档序（spTree 顺序）展开形状，组合形状完全展开、**无深度上限**。

    ⚠ 判组**不能**读 `shape.shape_type`：python-pptx 对没有可识别几何的 `p:sp` /
    `p:contentPart` 在这个 property 里 **raise NotImplementedError**（getattr 的
    default 挡不住抛出），一个畸形形状就会让整份 deck 掉回 parse_pptx_basic、表格
    全丢。GroupShape 是唯一带 `.shapes` 集合的形状类型，按它判既准又不触发计算。

    迭代显式栈而非递归：形状树是有限 XML 树、无环，任何深度截断都是对用户数据的
    静默丢弃（codex R3 P2——超深嵌套组会整支消失），栈遍历天然免疫 Python 递归
    上限，也就不需要「护栏」深度常量。
    """
    pending = list(shapes)
    pending.reverse()
    while pending:
        shape = pending.pop()
        if hasattr(shape, "shapes"):
            children = list(shape.shapes)
            children.reverse()
            pending.extend(children)
            continue
        yield shape


def _pptx_table_html(table: Any) -> str:
    """把 python-pptx 的表格重建成最小 HTML，与 MinerU 的 table_html 口径一致。

    合并单元格走 python-pptx 1.0 的公开 API：被覆盖的格（`is_spanned`）跳过，起始格
    （`is_merge_origin`）写 rowspan/colspan——否则一个跨 3 列的表头会被拍成 3 个格、
    其中 2 个空。
    """
    rows: List[str] = []
    for row in table.rows:
        cells: List[str] = []
        for cell in row.cells:
            if getattr(cell, "is_spanned", False):
                continue
            attrs = ""
            if getattr(cell, "is_merge_origin", False):
                span_height = int(getattr(cell, "span_height", 1) or 1)
                span_width = int(getattr(cell, "span_width", 1) or 1)
                if span_height > 1:
                    attrs += f' rowspan="{span_height}"'
                if span_width > 1:
                    attrs += f' colspan="{span_width}"'
            cells.append(f"<td{attrs}>{html_module.escape(' '.join(cell.text.split()))}</td>")
        rows.append(f"<tr>{''.join(cells)}</tr>")
    return f"<table>{''.join(rows)}</table>" if rows else ""


def parse_pptx_pythonpptx(source_id: str, path: Path) -> List[SourceElement]:
    """Structured PPTX fallback: python-pptx shapes, tables, charts and notes.

    The raw-XML extractor (parse_pptx_basic) only walks `p:sp` shapes, so slide
    tables (`p:graphicFrame` / `a:tbl`) and chart titles were dropped entirely.
    python-pptx exposes them through the public shape API.

    Pictures are skipped on purpose: this parser never persists image assets
    (the file may be a short-lived downloaded temp file), and python-pptx 1.0
    has no public alt-text accessor worth reaching into private XML for.

    Deliberately fail-open: missing wheel, unreadable package or zero elements
    falls through to parse_pptx_basic.
    """
    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
        elements: List[SourceElement] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            shape_index = 0
            table_index = 0
            for shape in _pptx_iter_shapes(slide.shapes):
                # per-shape 兜底：python-pptx 对畸形/未实现的形状会在属性访问处抛
                # 异常，坏形状只丢自己，不连坐同页其他形状、更不连坐整份 deck。
                # 这一层仍嵌在整函数的 fail-open except 之内——语义是「局部宽容、
                # 整体兜底」。
                try:
                    if getattr(shape, "has_table", False):
                        html = _pptx_table_html(shape.table)
                        parts = _split_table_html(html)
                        table_index += 1
                        for part_index, part in enumerate(parts, start=1):
                            text = _html_table_to_text(part)
                            if not text:
                                continue
                            label = f"PPTX slide {slide_index} table {table_index}"
                            metadata: Dict[str, Any] = {
                                "parser": "python-pptx",
                                "slide_number": slide_index,
                                "table_index": table_index,
                                "table_html": part,
                            }
                            if len(parts) > 1:
                                label = f"{label} part {part_index}"
                                metadata["table_part"] = part_index
                                metadata["table_parts"] = len(parts)
                            elements.append(
                                _element(source_id, "table", label, text, metadata)
                            )
                        continue
                    if getattr(shape, "has_chart", False):
                        chart = shape.chart
                        title = ""
                        if getattr(chart, "has_title", False):
                            title = " ".join(chart.chart_title.text_frame.text.split())
                        if not title:
                            continue
                        shape_index += 1
                        elements.append(
                            _element(
                                source_id,
                                "slide_text",
                                f"PPTX slide {slide_index} shape {shape_index}",
                                title,
                                {
                                    "parser": "python-pptx",
                                    "slide_number": slide_index,
                                    "shape_index": shape_index,
                                    "kind": "chart_title",
                                },
                            )
                        )
                        continue
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    text = " ".join(shape.text_frame.text.split())
                    if not text:
                        continue
                    shape_index += 1
                    elements.append(
                        _element(
                            source_id,
                            "slide_text",
                            f"PPTX slide {slide_index} shape {shape_index}",
                            text,
                            {
                                "parser": "python-pptx",
                                "slide_number": slide_index,
                                "shape_index": shape_index,
                            },
                        )
                    )
                except Exception:
                    continue

            if not slide.has_notes_slide:
                continue
            # notes_text_frame 可以是 None（备注母版没有 body 占位符，或占位符被删
            # 而 notesSlide part 仍在），官方文档写明了这一点。
            notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            if notes_frame is None:
                continue
            notes = " ".join(notes_frame.text.split())
            # Notes placeholders often echo the bare slide number; skip those.
            if not notes or notes == str(slide_index):
                continue
            elements.append(
                _element(
                    source_id,
                    "speaker_notes",
                    f"PPTX slide {slide_index} notes",
                    notes,
                    {"parser": "python-pptx", "slide_number": slide_index, "kind": "notes"},
                )
            )
        if elements:
            return elements
    except Exception:
        # Missing wheel / malformed package: one final try with the raw-XML
        # extractor so ingestion still completes.
        pass
    return parse_pptx_basic(source_id, path)


def parse_pptx_basic(source_id: str, path: Path) -> List[SourceElement]:
    elements: List[SourceElement] = []
    namespace = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }

    def shape_text(node: ElementTree.Element) -> str:
        values = [
            run.text.strip()
            for run in node.findall(".//a:t", namespace)
            if run.text and run.text.strip()
        ]
        return " ".join(values)

    with zipfile.ZipFile(path) as archive:
        slide_names = sorted(
            name
            for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
        for slide_index, slide_name in enumerate(slide_names, start=1):
            root = ElementTree.fromstring(archive.read(slide_name))
            shapes = root.findall(".//p:sp", namespace)
            shape_index = 0
            for shape in shapes:
                text = shape_text(shape)
                if not text:
                    continue
                shape_index += 1
                elements.append(
                    _element(
                        source_id,
                        "slide_text",
                        f"PPTX slide {slide_index} shape {shape_index}",
                        text,
                        {
                            "parser": "pptx",
                            "slide_number": slide_index,
                            "shape_index": shape_index,
                        },
                    )
                )
            if shape_index == 0:
                # No shape-level text (e.g. minimal XML); fall back to whole slide.
                text = shape_text(root)
                if text:
                    elements.append(
                        _element(
                            source_id,
                            "slide_text",
                            f"PPTX slide {slide_index}",
                            text,
                            {"parser": "pptx", "slide_number": slide_index},
                        )
                    )

        notes_names = sorted(
            name
            for name in archive.namelist()
            if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")
        )
        for notes_name in notes_names:
            match = re.search(r"notesSlide(\d+)\.xml", notes_name)
            notes_number = int(match.group(1)) if match else 0
            root = ElementTree.fromstring(archive.read(notes_name))
            text = shape_text(root).strip()
            # Notes placeholders often echo the bare slide number; skip those.
            if not text or text == str(notes_number):
                continue
            elements.append(
                _element(
                    source_id,
                    "speaker_notes",
                    f"PPTX slide {notes_number} notes",
                    text,
                    {"parser": "pptx", "slide_number": notes_number, "kind": "notes"},
                )
            )
    return elements


def parse_pdf(
    source_id: str,
    path: Path,
    file_name: str = "",
    mineru_client: Any = None,
    persist_image: Any = None,
) -> List[SourceElement]:
    """Compatibility wrapper for the provider-free Python PDF fallback."""
    return parse_pdf_python(source_id, path)


def parse_pdf_python(source_id: str, path: Path) -> List[SourceElement]:
    """High-quality local PDF fallback: PyMuPDF4LLM, then pypdf as last resort.

    Page-chunked Markdown preserves headings, multi-column reading order and
    tables substantially better than plain text extraction. Image files are
    intentionally not emitted here: the source may be a short-lived downloaded
    URL temp file, and dangling image paths must never enter persisted elements.
    """
    try:
        import pymupdf4llm

        page_chunks = pymupdf4llm.to_markdown(
            str(path),
            page_chunks=True,
            show_progress=False,
            write_images=False,
            embed_images=False,
            table_output="html",
        )
        elements: List[SourceElement] = []
        if isinstance(page_chunks, list):
            for page_number, chunk in enumerate(page_chunks, start=1):
                if not isinstance(chunk, dict):
                    continue
                markdown = str(chunk.get("text", "") or "")
                if not markdown.strip():
                    continue
                for ordinal, element in enumerate(
                    parse_markdown_text(source_id, markdown), start=1
                ):
                    metadata = dict(element.metadata)
                    metadata.update(
                        {
                            "parser": "pymupdf4llm",
                            "page_number": page_number,
                            "source_format": "pdf",
                        }
                    )
                    elements.append(
                        element.model_copy(
                            update={
                                "location_label": (
                                    f"PDF p.{page_number} {element.element_type} {ordinal}"
                                ),
                                "metadata": metadata,
                            }
                        )
                    )
        if elements:
            return elements
    except Exception:
        # The fallback itself is deliberately fail-open. A missing wheel,
        # unsupported PDF feature, or OCR/layout error gets one final pypdf try.
        pass
    return parse_pdf_pypdf(source_id, path)


def parse_pdf_pypdf(source_id: str, path: Path) -> List[SourceElement]:
    """Last-resort offline PDF fallback when PyMuPDF4LLM is unavailable/fails.

    Uses pypdf's layout-aware extraction (better reading order and column/row
    spacing than the default mode) and segments each page into paragraph and
    heading elements instead of one flattened blob. No formula/table fidelity —
    that requires MinerU — but gives a usable structured baseline offline.
    """
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("PDF parser dependency pypdf is not installed") from exc

    reader = PdfReader(str(path))
    elements: List[SourceElement] = []
    for page_index, page in enumerate(reader.pages, start=1):
        text = _extract_pdf_page_text(page)
        if not text.strip():
            continue
        para_index = 0
        heading_index = 0
        for is_heading, body in _segment_pdf_blocks(text):
            if is_heading:
                heading_index += 1
                elements.append(
                    _element(
                        source_id,
                        "heading",
                        f"PDF p.{page_index} heading {heading_index}",
                        body,
                        {"parser": "pdf", "page_number": page_index, "heading_index": heading_index},
                    )
                )
            else:
                para_index += 1
                elements.append(
                    _element(
                        source_id,
                        "page_text",
                        f"PDF p.{page_index} paragraph {para_index}",
                        body,
                        {"parser": "pdf", "page_number": page_index, "paragraph_index": para_index},
                    )
                )
    return elements


def _extract_pdf_page_text(page: Any) -> str:
    """Prefer pypdf layout mode; fall back to default extraction."""
    try:
        layout = page.extract_text(extraction_mode="layout") or ""
        if layout.strip():
            return layout
    except Exception:
        pass
    return page.extract_text() or ""


def _segment_pdf_blocks(text: str) -> List[tuple[bool, str]]:
    """Split page text into (is_heading, body) blocks on blank lines.

    Lines within a block are joined (wrapped text), intra-line runs of spaces
    are collapsed. A single short line without terminal punctuation is treated
    as a heading so document structure survives into extraction/citations.
    """
    blocks: List[tuple[bool, str]] = []
    for chunk in re.split(r"\n\s*\n", text):
        lines = [" ".join(line.split()) for line in chunk.splitlines() if line.strip()]
        if not lines:
            continue
        body = " ".join(lines).strip()
        if not body:
            continue
        is_heading = (
            len(lines) == 1
            and len(body) <= 80
            and not body.endswith((".", "。", ";", "；", ":", "：", ",", "，"))
        )
        blocks.append((is_heading, body))
    return blocks


def mineru_content_list_to_elements(
    source_id: str,
    content_list: List[dict],
    label_prefix: str = "PDF",
    images: Dict[str, bytes] | None = None,
    persist_image: Any = None,
) -> List[SourceElement]:
    """Map MinerU's content_list blocks into structured SourceElements.

    Formulas are kept as LaTeX (element_type "formula"), tables keep their HTML
    in metadata while exposing a flattened text body, and headings keep their
    level. page_idx is 0-based in MinerU; we display 1-based page numbers.
    """
    elements: List[SourceElement] = []
    counters: Dict[int, int] = {}
    for block in content_list:
        if not isinstance(block, dict):
            continue
        page = int(block.get("page_idx", 0)) + 1
        counters[page] = counters.get(page, 0) + 1
        ordinal = counters[page]
        block_type = str(block.get("type", "")).lower()

        if block_type in {"text", "title"}:
            text = str(block.get("text", "")).strip()
            if not text:
                continue
            level = int(block.get("text_level", 0) or 0)
            element_type = "heading" if level >= 1 else "paragraph"
            elements.append(
                _element(
                    source_id,
                    element_type,
                    f"{label_prefix} p.{page} block {ordinal}",
                    text,
                    {"parser": "mineru", "page_number": page, "text_level": level, "source_format": label_prefix.lower()},
                )
            )
        elif block_type == "equation":
            latex = _strip_math_delimiters(str(block.get("text", "")))
            if not latex:
                continue
            elements.append(
                _element(
                    source_id,
                    "formula",
                    f"{label_prefix} p.{page} formula {ordinal}",
                    latex,
                    {
                        "parser": "mineru",
                        "page_number": page,
                        "text_format": str(block.get("text_format", "latex")),
                        "source_format": label_prefix.lower(),
                    },
                )
            )
        elif block_type == "table":
            html = str(block.get("table_body", ""))
            caption = " ".join(_as_list(block.get("table_caption")))
            body_text = _html_table_to_text(html)
            text = " ".join(part for part in (caption, body_text) if part).strip()
            if not text:
                continue
            elements.append(
                _element(
                    source_id,
                    "table",
                    f"{label_prefix} p.{page} table {ordinal}",
                    text,
                    {
                        "parser": "mineru",
                        "page_number": page,
                        "table_html": html,
                        "caption": caption,
                        "source_format": label_prefix.lower(),
                    },
                )
            )
        elif block_type == "image":
            img_path = str(block.get("img_path", "")).strip()
            img_name = Path(img_path).name
            caption = " ".join(
                _as_list(block.get("image_caption")) + _as_list(block.get("image_footnote"))
            ).strip()
            asset_id = ""
            if img_name and images and persist_image and img_name in images:
                try:
                    asset_id = persist_image(images[img_name], img_name) or ""
                except Exception:
                    asset_id = ""  # 抽图/写盘失败绝不影响文本产出
            text = caption or f"{label_prefix} p.{page} 图 {ordinal}"
            metadata: Dict[str, Any] = {
                "parser": "mineru",
                "page_number": page,
                "source_format": label_prefix.lower(),
            }
            if asset_id:
                metadata["asset_id"] = asset_id
            if caption:
                metadata["caption"] = caption
            elements.append(
                _element(
                    source_id,
                    "image",
                    f"{label_prefix} p.{page} image {ordinal}",
                    text,
                    metadata,
                )
            )
        else:
            # Lists and any other block types: keep whatever text is present.
            text = str(block.get("text", "")).strip()
            if not text:
                continue
            elements.append(
                _element(
                    source_id,
                    "paragraph",
                    f"{label_prefix} p.{page} block {ordinal}",
                    text,
                    {"parser": "mineru", "page_number": page, "block_type": block_type, "source_format": label_prefix.lower()},
                )
            )
    return elements


def _as_list(value: Any) -> List[str]:
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()]
    if value:
        return [str(value).strip()]
    return []


def _strip_math_delimiters(text: str) -> str:
    cleaned = text.strip()
    for marker in ("$$", "\\[", "\\]", "$"):
        cleaned = cleaned.replace(marker, " ")
    return " ".join(cleaned.split())


def _top_level_table_rows(raw: str) -> List[str]:
    """取出一段 <table> 原始 HTML 里**顶层**的 <tr>…</tr> 片段。

    只认顶层：Word 允许表格嵌套，嵌套表的 <tr> 不是切点。深度按 <table>/</table>
    计数，嵌套表整体落在它所属的那个顶层 <tr> 里。
    """
    rows: List[str] = []
    depth = 0
    row_start = -1
    for match in re.finditer(r"(?i)<(/?)(table|tr)\b[^>]*>", raw):
        closing, tag = match.group(1), match.group(2).lower()
        if tag == "table":
            depth += 1 if not closing else -1
            continue
        if depth != 1:
            continue
        if not closing:
            if row_start < 0:
                row_start = match.start()
        elif row_start >= 0:
            rows.append(raw[row_start : match.end()])
            row_start = -1
    return rows


def _row_max_rowspan(row: str) -> int:
    """一行里最大的 rowspan 值（无 rowspan 记 1）。只认 <td>/<th> 上的属性。"""
    spans = [
        int(match.group(1))
        for match in re.finditer(
            r"(?i)<t[dh]\b[^>]*\browspan\s*=\s*[\"']?(\d+)", row
        )
    ]
    return max(spans, default=1)


def _split_table_html(raw: str, rows_per_element: int = _FALLBACK_TABLE_ROWS_PER_ELEMENT) -> List[str]:
    """按 <tr> 边界把超长表格 HTML 切成多段，每段仍是良构的 <table>。

    行数不超限时**原样返回**（单段、逐字不变），普通表格零行为变化。需要切分时每段
    重新用原始的 <table …> 开标签包裹：thead/tbody/caption 这类分组包装被丢弃，但
    表头语义真正所在的 <th> 单元格在行里、且 thead 的行本来就排在最前，因此自然归
    属首段。

    切点只落在**没有竖向合并跨越**的行边界上（codex R2 P2）：rowspan 跨段会让后段
    整列缺失、其余单元格全部错位。有活跃 rowspan 时顺延收段，段长因此可能超过
    rows_per_element；合并一路连到表尾的退化形态就整表一段——那等于该表不切，
    与「行数不超限原样返回」同一语义，宁可段长也不产出结构错误的碎片。
    """
    if not raw:
        return []
    rows = _top_level_table_rows(raw)
    if len(rows) <= rows_per_element:
        return [raw]
    open_match = re.search(r"(?i)<table\b[^>]*>", raw)
    open_tag = open_match.group(0) if open_match else "<table>"
    parts: List[str] = []
    part_start = 0
    covered_ahead = 0  # 当前行之后仍被上方 rowspan 覆盖的行数
    for index, row in enumerate(rows):
        covered_ahead = max(covered_ahead - 1, _row_max_rowspan(row) - 1)
        filled = index - part_start + 1
        if filled >= rows_per_element and covered_ahead == 0:
            parts.append(f"{open_tag}{''.join(rows[part_start : index + 1])}</table>")
            part_start = index + 1
    if part_start < len(rows):
        parts.append(f"{open_tag}{''.join(rows[part_start:])}</table>")
    return parts


def _html_table_to_text(html: str) -> str:
    """把表格 HTML 压成一行可检索的纯文本（`|` 分列、`;` 分行）。

    实体反转义必须发生在**剥标签之后**：`&lt;b&gt;` 先反转义就会变成真标签被下一步
    吃掉。这条同时修好了 MinerU 表格路径的同款既有问题——`table_html` 里的 `R&amp;D`
    此前原样进 embedding/FTS/展示，属净修复。`table_html` metadata 保持转义不变。
    """
    if not html:
        return ""
    text = re.sub(r"(?i)</t[dh]>", " | ", html)
    text = re.sub(r"(?i)</tr>", "\n", text)
    text = re.sub(r"<[^>]+>", " ", text)
    text = html_module.unescape(text)
    rows = [" ".join(row.split()).strip(" |") for row in text.splitlines()]
    return " ; ".join(row for row in rows if row)


def _element(
    source_id: str,
    element_type: str,
    location_label: str,
    text: str,
    metadata: Dict[str, Any],
) -> SourceElement:
    # 代码块/表格保真（保留换行/结构）；其余压平空白。
    if element_type in ("code_block", "table"):
        clean_text = text.strip("\n")
    else:
        clean_text = " ".join(text.split())
    return SourceElement(
        id="",
        source_id=source_id,
        element_type=element_type,
        location_label=location_label,
        text=clean_text,
        metadata=metadata,
    )
